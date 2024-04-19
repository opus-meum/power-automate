"""
Created on Mon Nov 13 2023

@author: ehlke_hepworth
"""
#company = 'EatYourGreens'

#def capability_report(company):

import numpy as np
import matplotlib.pyplot as plt
import matplotlib.font_manager as fm
import pandas as pd
import os
from matplotlib.font_manager import FontProperties
import matplotlib.pyplot as plt

font_prop = FontProperties(fname='fonts/AvenirNextCyr-Thin.ttf')
plt.rcParams['font.family'] = 'Avenir Next'

data = pd.read_csv('Capability Assessment Survey.csv').iloc[-1]
company = data.iloc[7]


dir_ = 'results/'
df = pd.read_csv(str(dir_)+'Summary_'+str(company)+'_.csv')
print(df)

fig1, ax = plt.subplots(1, figsize =(17, 10))

##set up heptagon
charcoal_color = '#425369'

a0 = 1.5
a1 = 1.2
a2 = 0.9
a3 = 0.6
a4 = 0.3

x0 = [a0 * np.cos(i * 2 * np.pi / 7) for i in range(8)]
y0 = [a0 * np.sin(i * 2 * np.pi / 7) for i in range(8)]

x1 = [a1 * np.cos(i * 2 * np.pi / 7) for i in range(8)]
y1 = [a1 * np.sin(i * 2 * np.pi / 7) for i in range(8)]

x2 = [a2 * np.cos(i * 2 * np.pi / 7) for i in range(8)]
y2 = [a2 * np.sin(i * 2 * np.pi / 7) for i in range(8)]

x3 = [a3 * np.cos(i * 2 * np.pi / 7) for i in range(8)]
y3 = [a3 * np.sin(i * 2 * np.pi / 7) for i in range(8)]

x4 = [a4 * np.cos(i * 2 * np.pi / 7) for i in range(8)]
y4 = [a4 * np.sin(i * 2 * np.pi / 7) for i in range(8)]

x_w = [1.77, 1, -0.3337814009344715, -1.65, -1.6,\
          -0.3337814009344719, 1,1]
y_w = [0, 1.25, 1.53, 0.6508256086763373, -0.650825608676337,\
          -1.53, -1.23,0]

words = ['PROCESSES', 'TALENT','STRATEGY','TECHNOLOGY','REPORTING','MEASUREMENT','DATA']

corners = []
for i in range(8):
   x = [a * np.cos(i * 2 * np.pi / 7) for a in [a1, a2, a3, a4]]
   y = [a * np.sin(i * 2 * np.pi / 7) for a in [a1, a2, a3, a4]]
   corners.append(list(zip(x, y)))
   x_coordinates_1 = [1.2, 0.7481877622304803, -0.26702512074757717, -1.0811626414829028, -1.081162641482903, -0.2670251207475775, 0.74818776223048, 1.2]
   y_coordinates_1 = [0.0, 0.9381977789616357, 1.1699134946181884, 0.5206604869410698, -0.5206604869410696, -1.1699134946181884, -0.9381977789616358, -2.9391523179536476e-16]

print(df['Processes'])

if df['Processes'].iloc[0] == 5:
    x_pr = x0[0];
    y_pr = y0[0]
if 4 < df['Processes'].iloc[0] < 5:
    x_pr = (x0[0]+x1[0])/2;
    y_pr = (y0[0]+y1[0])/2
if df['Processes'].iloc[0] == 4:
    x_pr = x1[0];
    y_pr = y1[0]
if 3< df['Processes'].iloc[0] < 4:
    x_pr = (x1[0]+x2[0])/2;
    y_pr = (y1[0]+y2[0])/2
if df['Processes'].iloc[0] == 3:
    x_pr = x2[0];
    y_pr = y2[0]
if 2 < df['Processes'].iloc[0] <3:
    x_pr = (x2[0]+x3[0])/2;
    y_pr = (y2[0]+y3[0])/2
if df['Processes'].iloc[0] == 2:
    x_pr = x3[0];
    y_pr = y3[0]
if 1 < df['Processes'].iloc[0] <2 :
    x_pr = (x3[0]+x4[0])/2;
    y_pr = (y3[0]+y4[0])/2
if df['Processes'].iloc[0] ==1:
    x_pr = x4[0];
    y_pr = y4[0]
if 0 < df['Processes'].iloc[0] <1 :
    x_pr = (x4[0])/2;
    y_pr = (y4[0])/2
if df['Processes'].iloc[0] == 0:
    x_pr = 0;
    y_pr = 0
   
if df['Talent'].iloc[0] == 5:
    x_tal = x0[1];
    y_tal = y0[1]
if 4 < df['Talent'].iloc[0] < 5:
    x_tal = (x0[1]+x1[1])/2;
    y_tal = (y0[1]+y1[1])/2
if df['Talent'].iloc[0] == 4:
    x_tal = x1[1];
    y_tal = y1[1]
if 3 < df['Talent'].iloc[0] < 4:
    x_tal = (x1[1]+x2[1])/2;
    y_tal = (y1[1]+y2[1])/2
if df['Talent'].iloc[0] == 3:
    x_tal = x2[1];
    y_tal = y2[1]
if 2 < df['Talent'].iloc[0] < 3:
    x_tal = (x2[1]+x3[1])/2;
    y_tal = (y2[1]+y3[1])/2
if df['Talent'].iloc[0] == 2:
    x_tal = x3[1];
    y_tal = y3[1]
if 1 < df['Talent'].iloc[0] <2 :
    x_tal = (x3[1]+x4[1])/2;
    y_tal = (y3[1]+y4[1])/2
if df['Talent'].iloc[0] ==1:
    x_tal = x4[1];
    y_tal = y4[1]
if 0 < df['Talent'].iloc[0] <1 :
    x_tal = (x4[1])/2;
    y_tal = (y4[1])/2
if df['Talent'].iloc[0] == 0:
    x_tal = 0;
    y_tal = 0


if df['Strategy'].iloc[0] == 5:
    x_st = x0[2];
    y_st = y0[2]
if 4 < df['Strategy'].iloc[0] < 5:
    x_st = (x0[2]+x1[2])/2;
    y_st = (y0[2]+y1[2])/2
if df['Strategy'].iloc[0] == 4:
   x_st = x1[2];
   y_st = y1[2]
if 3 < df['Strategy'].iloc[0] < 4:
   x_st = (x1[2]+x2[2])/2;
   y_st = (y1[2]+y2[2])/2
if df['Strategy'].iloc[0] == 3:
   x_st = x2[2];
   y_st = y2[2]
if 2 < df['Strategy'].iloc[0] < 3:
   x_st = (x2[2]+x3[2])/2;
   y_st = (y2[2]+y3[2])/2
if df['Strategy'].iloc[0] == 2:
   x_st = x3[2];
   y_st = y3[2]
if 1 < df['Strategy'].iloc[0] <2 :
   x_st = (x3[2]+x4[2])/2;
   y_st = (y3[2]+y4[2])/2
if df['Strategy'].iloc[0] ==1:
   x_st = x4[2];
   y_st = y4[2]
if 0 < df['Strategy'].iloc[0] <1 :
   x_st = (x4[2])/2;
   y_st = (y4[2])/2
if df['Strategy'].iloc[0] == 0:
   x_st = 0;
   y_st = 0

if df['Technology'].iloc[0] == 5:
    x_te = x0[3];
    y_te = y0[3]
if 4 < df['Technology'].iloc[0] < 5:
    x_te = (x0[3]+x1[3])/2;
    y_te = (y0[3]+y1[3])/2
if df['Technology'].iloc[0] == 4:
   x_te = x1[3];
   y_te = y1[3]
if df['Technology'].iloc[0] > 3 and df['Technology'].iloc[0] < 4:
   x_te = (x1[3]+x2[3])/2;
   y_te = (y1[3]+y2[3])/2
if df['Technology'].iloc[0] == 3:
   x_te = x2[3];
   y_te = y2[3]
if 2 < df['Technology'].iloc[0] < 3:
   x_te = (x2[3]+x3[3])/2;
   y_te = (y2[3]+y3[3])/2
if df['Technology'].iloc[0] == 2:
   x_te = x3[3];
   y_te = y3[3]
if 1 < df['Technology'].iloc[0] <2 :
   x_te = (x3[3]+x4[3])/2;
   y_te = (y3[3]+y4[3])/2
if df['Technology'].iloc[0] ==1:
   x_te = x4[3];
   y_te = y4[3]
if 0 < df['Technology'].iloc[0] <1 :
   x_te = (x4[3])/2;
   y_te = (y4[3])/2
if df['Technology'].iloc[0] == 0:
   x_te = 0;
   y_te = 0

if df['Reporting'].iloc[0] == 5:
    x_re = x0[4];
    y_re = y0[4]
if 4 < df['Reporting'].iloc[0] < 5:
    x_re = (x0[4]+x1[4])/2;
    y_re = (y0[4]+y1[4])/2
if df['Reporting'].iloc[0] == 4:
   x_re = x1[4];
   y_re = y1[4]
if 3 < df['Reporting'].iloc[0] < 4:
   x_re = (x1[4]+x2[4])/2;
   y_re = (y1[4]+y2[4])/2
if df['Reporting'].iloc[0] == 3:
   x_re = x2[4];
   y_re = y2[4]
if 2 < df['Reporting'].iloc[0] < 3:
   x_re = (x2[4]+x3[4])/2;
   y_re = (y2[4]+y3[4])/2
if df['Reporting'].iloc[0] == 2:
   x_re = x3[4];
   y_re = y3[4]
if 1 < df['Reporting'].iloc[0] <2 :
   x_re = (x3[4]+x4[4])/2;
   y_re = (y3[4]+y4[4])/2
if df['Reporting'].iloc[0] == 1:
   x_re = x4[4];
   y_re = y4[4]
if 0 < df['Reporting'].iloc[0] <1 :
   x_re = (x4[4])/2;
   y_re = (y4[4])/2
if df['Reporting'].iloc[0] == 0:
   x_re = 0;
   y_re = 0
          
if df['Measurement'].iloc[0] == 5:
    x_me = x0[5];
    y_me = y0[5]
if 4 < df['Measurement'].iloc[0] < 5:
    x_me = (x0[5]+x1[5])/2;
    y_me = (y0[5]+y1[5])/2
if df['Measurement'].iloc[0] == 4:
   x_me = x1[5];
   y_me = y1[5]
if 3 < df['Measurement'].iloc[0] < 4:
   x_me = (x1[5]+x2[5])/2;
   y_me = (y1[5]+y2[5])/2
if df['Measurement'].iloc[0] == 3:
   x_me = x2[5];
   y_me = y2[5]
if 2 < df['Measurement'].iloc[0] < 3:
   x_me = (x2[5]+x3[5])/2;
   y_me = (y2[5]+y3[5])/2
if df['Measurement'].iloc[0] == 2:
   x_me = x3[5];
   y_me = y3[5]
if 1 < df['Measurement'].iloc[0] <2 :
   x_me = (x3[5]+x4[5])/2;
   y_me = (y3[5]+y4[5])/2
if df['Measurement'].iloc[0] == 1:
   x_me = x4[5];
   y_me = y4[5]
if 0 < df['Processes'].iloc[0] <1 :
   x_me = (x4[5])/2;
   y_me = (y4[5])/2
if df['Measurement'].iloc[0] == 0:
   x_me = 0;
   y_me = 0

if df['Data'].iloc[0] == 5:
    x_da = x0[6];
    y_da = y0[6]
if 4 < df['Data'].iloc[0] < 5:
    x_da = (x0[6]+x1[6])/2;
    y_da = (y0[6]+y1[6])/2
if df['Data'].iloc[0] == 4:
   x_da = x1[6];
   y_da = y1[6]
if 3 < df['Data'].iloc[0] < 4:
   x_da = (x1[6]+x2[6])/2;
   y_da = (y1[6]+y2[6])/2
if df['Data'].iloc[0] == 3:
   x_da = x2[6];
   y_da = y2[6]
if 2 < df['Data'].iloc[0] < 3:
   x_da = (x2[6]+x3[6])/2;
   y_da = (y2[6]+y3[6])/2
if df['Data'].iloc[0] == 2:
   x_da = x3[6];
   y_da = y3[6]
if 1 < df['Data'].iloc[0] <2 :
   x_da = (x3[6]+x4[6])/2;
   y_da = (y3[6]+y4[6])/2
if df['Data'].iloc[0] == 1:
   x_da = x4[6];
   y_da = y4[6]
if 0 < df['Processes'].iloc[0] <1 :
   x_da = (x4[6])/2;
   y_da = (y4[6])/2
if df['Data'].iloc[0] == 0:
   x_da = 0;
   y_da = 0
# Print the coordinates for all corners of the heptagon

####set up current level of company
x_current = [x_pr,x_tal,x_st,x_te,x_re,x_me,x_da,x_pr]
y_current = [y_pr,y_tal,y_st,y_te,y_re,y_me,y_da,y_pr]

print(x_current)
print(y_current)
#x_future  = [0.901,0.568, -.226, -.535,-.654,-.160,.613,.901]
#y_future  = [0.0009,0.693,1.010,.249,-.330,-.731 ,-.772,.0009]

x = [1.5, 0.9352347027881004, -0.3337814009344715, -1.3514533018536286, -1.3514533018536288,\
      -0.3337814009344719, 0.9352347027881001,1.5]
y = [0, 1.1727472237020446, 1.4623918682727355, 0.6508256086763373, -0.650825608676337,\
      -1.4623918682727355, -1.1727472237020449,0]
#font_path = fm.findfont(fm.FontProperties(family='Avenir Next light'))

# Plot the heptagon
plt.plot(x,y, color='grey', linewidth=5, label = "Mature State")
plt.plot(x1,y1, color='grey',linewidth=.5)
plt.plot(x2,y2, color='grey',linewidth=.5)
plt.plot(x3,y3, color='grey',linewidth=.5)
plt.plot(x4,y4, color='grey',linewidth=.5, label ="Maturity Steps")


plt.plot(x_current,y_current,color=charcoal_color,linewidth=5, label = "Current State")
#   plt.plot(x_future,y_future,color='red',linewidth=5, label="Future State")

for i, word in enumerate(words):
   plt.text(x_w[i], y_w[i], word, ha='center', va='center', fontname='Avenir Next', fontweight='light',color=charcoal_color,size=12)
#legend = plt.legend(prop={'family':'Avenir Next'})
legend = plt.legend()
for text in legend.get_texts():
   text.set_color('#425369')
plt.xticks([])
plt.yticks([])
# Set the aspect ratio to equal
plt.axis('equal')

# Find the minimum value across all dimensions
min_value = df.min(numeric_only=True).min()
# Find all dimensions where the minimum value is found
min_dimensions = df.min(numeric_only=True)[df.min(numeric_only=True) == min_value].index.tolist()

if len(min_dimensions) > 1:
    min_dimensions_str = ', '.join(min_dimensions[:-1]) + ' and ' + min_dimensions[-1]
else:
    min_dimensions_str = min_dimensions[0]
          
plt.annotate(f"This figure highlights the greatest gaps, \
            \nrepresenting the degree of shift required, \
            \nbetween the average current state (depicted \
            \nby the thick, dark line) and the potential \
            \nfuture state (thick,grey line) across \
            \n{min_dimensions_str}.",
         xy=(1, 0),  # point to label
         xytext=(3, 0),  # location of text
         ha='left', va='center',
#         arrowprops=dict(facecolor='black', shrink=0.05),
#         bbox=dict(boxstyle="round,pad=0.3", fc="whitesmoke", alpha=0.8),
         fontsize=20,
         fontname='Avenir Next', color='#425369'
         )


# Show the plot
plt.savefig(str(dir_)+str(company)+'heptagon_plot.png',dpi=300, bbox_inches='tight')

#%% 
##------------------------------------------------
import pandas as pd
import matplotlib.pyplot as plt
#import seaborn as sns 
import numpy as np 


#%%
fig2,ax = plt.subplots(1, figsize =(20, 9))

data = pd.read_csv(str(dir_)+'Strategy_'+str(company)+'_.csv',sep=',')
strategy = data['max'][:-1]
dim_st   = data['Dimension'][:-1]
data = pd.read_csv(str(dir_)+'Talent_'+str(company)+'_.csv',sep=',')
talent  = data['max'][:-1]
dim_ta  = data['Dimension'][:-1]
data = pd.read_csv(str(dir_)+'Processes_'+str(company)+'_.csv',sep=',')
processes  = data['max'][:-1]
dim_pr  = data['Dimension'][:-1]
data = pd.read_csv(str(dir_)+'Data_'+str(company)+'_.csv',sep=',')
data_  = data['max'][:-1]
dim_da  = data['Dimension'][:-1]
data = pd.read_csv(str(dir_)+'Technology_'+str(company)+'_.csv',sep=',')
tech  = data['max'][:-1]
dim_te  = data['Dimension'][:-1]
data = pd.read_csv(str(dir_)+'Reporting_'+str(company)+'_.csv',sep=',')
reporting  = data['max'][:-1]
dim_re  = data['Dimension'][:-1]
data = pd.read_csv(str(dir_)+'Measurement_'+str(company)+'_.csv',sep=',')
measurement  = data['max'][:-1]
dim_me  = data['Dimension'][:-1]


#%%

##x- and y- axes
x_ax = (['STRATEGY', 'TALENT', 'PROCESSES', 'DATA', 'MEASUREMENT', 'REPORTING','TECHNOLOGY'])
y_ax = (['5','4','3','2','1','0'])

ax.set_xticks([0,1, 2, 3, 4, 5, 6])

ax.set_xticklabels(x_ax, position=(0,-0.01), font='Avenir Next',color='#425369')
#ax.set_xlabel('Capability Assessment Dimension',fontsize=25,\
#                        font='Avenir Next',color='#425369')
ax.xaxis.labelpad = 20
ax.tick_params(axis='x', labelsize=17)
ax.set_xlim(-0.5, 6.5)


ax.set_yticks([1,2,3,4,5,6])
ax.set_yticklabels(y_ax,color='#425369')
ax.tick_params(axis='y', labelsize=15.5)
# ax.set_ylabel('Steps Between Current State and Future State',fontsize=16,\
#           font='Avenir Next',color='#425369')
ax.yaxis.labelpad = 18
ax.set_ylim(0.5, 6.5)

c_s = ['#425369']
c_ta= ['#425369']
c_p = ['#425369']
c_d = ['#425369']
c_m = ['#425369']
c_r = ['#425369']
c_te= ['#425369']


##bars and lines
plt.bar(x_ax, [6,6,6,6,6,6,6], align='center', alpha=0.3,color='grey')
plt.grid(False)
st_wor = []
plt.hlines(6, 0-0.4, 0+0.4,linewidth=45, color='darkgrey')
plt.hlines(6, 1-0.4, 1+0.4,linewidth=45, color='darkgrey')
plt.hlines(6, 2-0.4, 2+0.4,linewidth=45, color='darkgrey')
plt.hlines(6, 3-0.4, 3+0.4,linewidth=45, color='darkgrey')
plt.hlines(6, 4-0.4, 4+0.4,linewidth=45, color='darkgrey')
plt.hlines(6, 5-0.4, 5+0.4,linewidth=45, color='darkgrey')
plt.hlines(6, 6-0.4, 6+0.4,linewidth=45, color='darkgrey')

    
for i in range(len(dim_st)):
    if strategy[i] > 3:
        linewidth = 75
    else:
        linewidth = 65
    plt.hlines(strategy[i]+1, 0-0.4, 0+0.4,linewidth=linewidth, color=c_s)
# plt.text(-0.01 * (i+1), strategy[i]+1.3-0.1*i, dim_st.iloc[i], ha='center', va='top', color='whitesmoke', fontsize=9,fontweight='bold')
#     plt.text(-0.4,5-strategy[i],dim_st[i],fontsize=10,font='Avenir Next',color='whitesmoke')
for i in range(len(dim_ta)):
    if talent[i] > 3:
        linewidth = 75
    else:
        linewidth = 65
    plt.hlines(talent[i]+1, 1-0.4, 1+0.4,linewidth=linewidth, color=c_ta)
   #plt.text(0.99+0 * (i+1), talent[i]+1.3-0.15*i, dim_ta.iloc[i], ha='center', va='top', color='whitesmoke', fontsize=9,fontweight='bold')

#       plt.text(0.7+0.1 * (i+1), talent[i]+1, str(i+1), ha='center', va='bottom', color='whitesmoke', fontsize=15)
#    plt.text(0.6,5-talent[i],dim_ta[i],fontsize=10,font='Avenir Next',color='whitesmoke')
for i in range(len(dim_pr)):
    if processes[i] > 3:
        linewidth = 75
    else:
        linewidth = 65
    plt.hlines(processes[i]+1, 2.0-0.4, 2.0+0.4,linewidth=linewidth, color=c_p)
# plt.text(1.99+0 * (i+1), processes[i]+1.3-0.15*i, dim_pr.iloc[i], ha='center', va='top', color='whitesmoke', fontsize=9,fontweight='bold')

#       plt.text(1.7+0.1 * (i+1), processes[i]+1, str(i+1), ha='center', va='bottom', color='whitesmoke', fontsize=15)
#    plt.text(1.6,5-processes[i],dim_pr[i],fontsize=10,font='Avenir Next',color='whitesmoke')   
for i in range(len(dim_da)):
    if data_[i] > 3:
        linewidth = 75
    else:
        linewidth = 65
    plt.hlines(data_[i]+1, 3-0.4, 3+0.4,linewidth=linewidth, color=c_d)
# plt.text(2.99+0 * (i+1), data_[i]+1.3-0.15*i, dim_da.iloc[i], ha='center', va='top', color='whitesmoke', fontsize=9,fontweight='bold')

#      plt.text(2.7+0.1 * (i+1), data_[i]+1, str(i+1), ha='center', va='bottom', color='whitesmoke', fontsize=15)
#    plt.text(2.6,5-data_[i],dim_da[i],fontsize=10,font='Avenir Next',color='whitesmoke')   
for i in range(len(dim_st)):
    if measurement[i] > 3:
        linewidth = 75
    else:
        linewidth = 65
    plt.hlines(measurement[i]+1, 4-0.4, 4+0.4, linewidth=linewidth, color=c_s)
##for i in range(len(dim_me)):
##   plt.hlines(measurement[i]+1, 4-0.4, 4+0.4,linewidth=75, color=c_m)
# plt.text(3.99, measurement[i]+1.3-0.15*i, dim_me.iloc[i], ha='center', va='top', color='whitesmoke', fontsize=9, fontweight='bold')
#       plt.text(3.7+0.1 * (i+1), measurement[i]+1, str(i+1), ha='center', va='bottom', color='whitesmoke', fontsize=15)
#    plt.text(3.6,5-measurement[i],dim_me[i],fontsize=10,font='Avenir Next',color='whitesmoke')   
for i in range(len(dim_re)):
    if reporting[i] > 3:
        linewidth = 75
    else:
        linewidth = 65
    plt.hlines(reporting[i]+1, 5-0.4, 5+0.4,linewidth=linewidth, color=c_r)
# plt.text(4.99+0 * (i+1), reporting[i]+1.3-0.15*i, dim_re.iloc[i], ha='center', va='top', color='whitesmoke', fontsize=9,fontweight='bold')

#       plt.text(4.7+0.1 * (i+1), reporting[i]+1, str(i+1), ha='center', va='bottom', color='whitesmoke', fontsize=15)
#    plt.text(4.6,5-reporting[i],dim_re[i],fontsize=10,font='Avenir Next',color='whitesmoke')   
for i in range(len(dim_te)):
    if tech[i] > 3:
        linewidth = 75
    else:
        linewidth = 65
    plt.hlines(tech[i]+1, 6-0.4, 6+0.4,linewidth=linewidth, color=c_te)
# plt.text(5.99+0 * (i+1), tech[i]+1.3-0.15*i, dim_te.iloc[i], ha='center', va='top', color='whitesmoke', fontsize=9, fontweight='bold')


# Adjusting text positioning for overlapping texts on hlines
def adjust_text_position(text_positions, current_position, direction=-1):
   """Adjust text position to avoid overlap, moving up or down."""
   adjusted_position = current_position
   while adjusted_position in text_positions:
      adjusted_position += 0.19 * direction  # Adjust by 0.15 units up or down
   text_positions.add(adjusted_position)
   return adjusted_position

# Strategy dimension text positioning
strategy_text_positions = set()
for i in range(len(dim_st)):
   if strategy[i] > 3:
      y_position = 1.4
   else:
      y_position = 1.3
   text_y = strategy[i] + y_position  # Initial Y position for text
   text_y = adjust_text_position(strategy_text_positions, text_y)
   plt.text(0, text_y, dim_st.iloc[i], ha='center', va='top', color='whitesmoke', fontsize=9, fontweight='bold')

# Talent dimension text positioning
talent_text_positions = set()
for i in range(len(dim_ta)):
   if talent[i] > 3:
      y_position = 1.4
   else:
      y_position = 1.3
   text_y = talent[i] + y_position # Initial Y position for text
   text_y = adjust_text_position(talent_text_positions, text_y)
   plt.text(1, text_y, dim_ta.iloc[i], ha='center', va='top', color='whitesmoke', fontsize=9, fontweight='bold')

# Processes dimension text positioning
processes_text_positions = set()
for i in range(len(dim_pr)):
   if processes[i] > 3:
      y_position = 1.4
   else:
      y_position = 1.3
   text_y = processes[i] + y_position  # Initial Y position for text
   text_y = adjust_text_position(processes_text_positions, text_y)
   plt.text(2, text_y, dim_pr.iloc[i], ha='center', va='top', color='whitesmoke', fontsize=9, fontweight='bold')

# Data dimension text positioning
data_text_positions = set()
for i in range(len(dim_da)):
   if data_[i] > 3:
      y_position = 1.4
   else:
      y_position = 1.3
   text_y = data_[i] + y_position # Initial Y position for text
   text_y = adjust_text_position(data_text_positions, text_y)
   plt.text(3, text_y, dim_da.iloc[i], ha='center', va='top', color='whitesmoke', fontsize=9, fontweight='bold')

# Measurement dimension text positioning
measurement_text_positions = set()
for i in range(len(dim_me)):
   if measurement[i] > 3:
      y_position = 1.4
   else:
      y_position = 1.3
   text_y = measurement[i] + y_position  # Initial Y position for text
   text_y = adjust_text_position(measurement_text_positions, text_y)
   plt.text(4, text_y, dim_me.iloc[i], ha='center', va='top', color='whitesmoke', fontsize=9, fontweight='bold')

# Reporting dimension text positioning
reporting_text_positions = set()
for i in range(len(dim_re)):
   if reporting[i] > 3:
      y_position = 1.4
   else:
      y_position = 1.3
   text_y = reporting[i] + y_position  # Initial Y position for text
   text_y = adjust_text_position(reporting_text_positions, text_y)
   plt.text(5, text_y, dim_re.iloc[i], ha='center', va='top', color='whitesmoke', fontsize=9, fontweight='bold')

# Technology dimension text positioning
tech_text_positions = set()
for i in range(len(dim_te)):
   if tech[i] > 3:
      y_position = 1.4
   else:
      y_position = 1.3
   text_y = tech[i] + y_position # Initial Y position for text
   text_y = adjust_text_position(tech_text_positions, text_y)
   plt.text(6, text_y, dim_te.iloc[i], ha='center', va='top', color='whitesmoke', fontsize=9, fontweight='bold')

  #     plt.text(5.7+0.1 * (i+1), tech[i]+1, str(i+1), ha='center', va='bottom', color='whitesmoke', fontsize=15)
#    plt.text(5.6,5-tech[i],dim_te[i],fontsize=10,font='Avenir Next',color='whitesmoke')   
"""
# Strategy dimension
# Initialize a dictionary to keep track of the last y-coordinate used for each hline
# Strategy dimension
    last_hline_y = -1
    text_count = 0
    for i in range(len(dim_st)):
      current_hline_y = strategy[i]+1
      if current_hline_y != last_hline_y:
        last_hline_y = current_hline_y
        text_y = current_hline_y + 0.15
        text_count = 0
      else:
        text_count += 1
        text_y = current_hline_y + 0.15 - 0.15 * text_count
      plt.hlines(current_hline_y, 0-0.4, 0+0.4, linewidth=55, color=c_s)
      plt.text(-0.01 * (i+1), text_y, dim_st.iloc[i], ha='center', va='top', color='whitesmoke', fontsize=9, fontweight='bold')
# Talent dimension
    last_hline_y = -1
    last_text_y = -1
    for i in range(len(dim_ta)):
       current_hline_y = talent[i]+1
       if current_hline_y != last_hline_y:
          last_hline_y = current_hline_y
          text_y = current_hline_y + 0.15
          last_text_y = text_y
       else:
          text_y = last_text_y - 0.15
          last_text_y = text_y
          plt.hlines(current_hline_y, 1-0.4, 1+0.4, linewidth=55, color=c_ta)
          plt.text(1, text_y-0.15, dim_ta.iloc[i], ha='center', va='top', color='whitesmoke', fontsize=9, fontweight='bold')
"""

dimensions = [strategy, talent, processes, data_, measurement, reporting, tech]
dimension_names = [dim_st, dim_ta, dim_pr, dim_da, dim_me, dim_re, dim_te]

"""
related_words=[]
for dim_data, dim_name in zip(dimensions, dimension_names):
    for i, value in enumerate(dim_data):
        if value == 2:
            related_words.append(dim_name.iloc[i])
"""
# Flatten all dim_data lists into a single list
all_values = [value for dim_data in dimensions for value in dim_data]

# Find the minimum value across all dimensions
overall_min_value = min(all_values)

# Now, you can check against this overall minimum value
related_words = []
for dim_data, dim_name in zip(dimensions, dimension_names):
    for i, value in enumerate(dim_data):
        if value == overall_min_value:
            related_words.append(dim_name.iloc[i])

##add test to hlines
plt.text(-0.2,6,'Future State',fontsize=11,font='Avenir Next',color='whitesmoke')
plt.text(0.8,6,'Future State',fontsize=11,font='Avenir Next',color='whitesmoke')
plt.text(1.8,6,'Future State',fontsize=11,font='Avenir Next',color='whitesmoke')
plt.text(2.8,6,'Future State',fontsize=11,font='Avenir Next',color='whitesmoke')
plt.text(3.8,6,'Future State',fontsize=11,font='Avenir Next',color='whitesmoke')
plt.text(4.8,6,'Future State',fontsize=11,font='Avenir Next',color='whitesmoke')
plt.text(5.8,6,'Future State',fontsize=11,font='Avenir Next',color='whitesmoke')

x_limits = plt.xlim()
y_limits = plt.ylim()

print("X limits:", x_limits)
print("Y limits:", y_limits)

ax.grid(False)
plt.savefig(str(dir_)+'CapAss.png',dpi=300,bbox_inches='tight')


import matplotlib.backends.backend_pdf as pdf
##Individual plots\
"""
fig3,ax = plt.subplots(1, figsize =(17, 10))

data = pd.read_csv(dir_+'Strategy_'+company+'_.csv',sep=',')
    strategy = data['max'][:-1]
    dim_st   = data['Dimension'][:-1]

    print(strategy)
    print(dim_st)
    # Plotting the bar graph for Strategy dimension
    maturity_dict = {'nascent': 1, 'emerging': 2, 'expanding': 3, 'optimising': 4, 'mature': 5}

    charcoal_color = '#425369'
    font = {'family': 'Avenir Next'}
    plt.rc('font', **font, weight='light')

    ax.bar(dim_st, strategy, color=charcoal_color, edgecolor='black', linewidth=1, alpha=0.7)
    for num in strategy:
        for key, value in maturity_dict.items():
           for i, level in enumerate(strategy):
              ax.text(dim_st[i], level + 0.1, list(maturity_dict.keys())[list(maturity_dict.values()).index(level)], color=charcoal_color, fontweight='light', ha='center',fontsize=18)
    plt.xticks(fontsize=25, color=charcoal_color)
    plt.yticks(fontsize=25, color=charcoal_color)
    
    plt.ylim(0, 5)
    plt.tight_layout()
    # Modernizing the graph with seaborn style
    import seaborn as sns
    sns.set_style("whitegrid")
    sns.despine(left=True)

    plt.gca().spines['top'].set_visible(False)
    plt.gca().spines['right'].set_visible(False)
    plt.gca().spines['bottom'].set_color('#dddddd')
    plt.gca().spines['left'].set_color('#dddddd')
    plt.gca().xaxis.label.set_color('#333333')
    plt.gca().yaxis.label.set_color('#333333')
    plt.gca().tick_params(colors='#333333')
      
   # plt.show(bbox_inches='tight')
    plt.savefig('/Users/ehlke/Desktop/strategy.png',dpi=300, bbox_inches='tight')
    
    from pptx import Presentation
    from pptx.util import Inches
    prs = Presentation('/Users/ehlke/Desktop/Client X Capability Assessment Report Template.pptx')  

    # Select the slide where you want to add the image
    slide = prs.slides[11]  # slide_index is the index of the slide
    # Define the width, height, and coordinates for the image
    left = Inches(1.1)
    top = Inches(1.8)
    width = Inches(5.5)
    height = Inches(3.5)
    
    slide.shapes.add_picture('/Users/ehlke/Desktop/strategy.png', left, top, width, height)

    fig4,ax = plt.subplots(1, figsize =(17, 10))

    data = pd.read_csv(dir_+'Talent_'+company+'_.csv',sep=',')
    talent  = data['max'][:-1]
    dim_ta  = data['Dimension'][:-1]

    # Plotting the bar graph for Strategy dimension
    maturity_dict = {'nascent': 1, 'emerging': 2, 'expanding': 3, 'optimising': 4, 'mature': 5}

    charcoal_color = '#425369'
    font = {'family': 'Avenir Next'}
    plt.rc('font', **font, weight='light')

    ax.bar(dim_ta, talent, color=charcoal_color, edgecolor='black', linewidth=1, alpha=0.7)
    for num in talent:
        for key, value in maturity_dict.items():
           for i, level in enumerate(talent):
              ax.text(dim_ta[i], level + 0.1, list(maturity_dict.keys())[list(maturity_dict.values()).index(level)], color=charcoal_color, fontweight='light', ha='center',fontsize=18)
    plt.xticks(fontsize=25, color=charcoal_color)
    plt.yticks(fontsize=25, color=charcoal_color)
    
    plt.ylim(0, 5)
    plt.tight_layout()
    # Modernizing the graph with seaborn style
    import seaborn as sns
    sns.set_style("whitegrid")
    sns.despine(left=True)

    plt.gca().spines['top'].set_visible(False)
    plt.gca().spines['right'].set_visible(False)
    plt.gca().spines['bottom'].set_color('#dddddd')
    plt.gca().spines['left'].set_color('#dddddd')
    plt.gca().xaxis.label.set_color('#333333')
    plt.gca().yaxis.label.set_color('#333333')
    plt.gca().tick_params(colors='#333333')
      
   # plt.show(bbox_inches='tight')
    plt.savefig('/Users/ehlke/Desktop/talent.png',dpi=300, bbox_inches='tight')
    """


from pptx import Presentation
from pptx.util import Inches
import seaborn as sns
import matplotlib.pyplot as plt
import pandas as pd

def create_plot_and_add_to_ppt(category, company, dir_, prs, slide_index):
   fig, ax = plt.subplots(1, figsize =(17, 10))
   #plt.rcParams['font.sans-serif'] = 'Dejavu Sans'

   data = pd.read_csv(f'{dir_}{category}_{company}_.csv', sep=',')
   values = data['max'][:-1]
   dimensions = data['Dimension'][:-1]

   maturity_dict = {'Nascent': 1, 'Emerging': 2, 'Expanding': 3, 'Optimising': 4, 'Mature': 5}
   charcoal_color = '#425369'
   #font = {'family': 'Dejavu Sans'}
   #plt.rc('font', **font, weight='light')

   num_bars = len(dimensions)
   # Dynamically adjust bar width based on the number of bars
   bar_width = num_bars / 4.2

   ax.bar(dimensions, values, width=bar_width, color=charcoal_color, edgecolor='black', linewidth=1, alpha=0.7)
   for num in values:
      for key, value in maturity_dict.items():
         for i, level in enumerate(values):
               ax.text(dimensions[i], level + 0.1, list(maturity_dict.keys())[list(maturity_dict.values()).index(level)],color=charcoal_color, fontweight='light', ha='center',fontsize=18)
   plt.xticks(color=charcoal_color, font='Avenir Next')
   plt.yticks(color=charcoal_color, font='Avenir Next')
   ax.tick_params(axis='both', which='major', labelsize=21)
   ax.tick_params(axis='x', colors=charcoal_color)

   plt.ylim(0, 5)
   plt.tight_layout()
   # Adjust x-axis limits to better display a single bar
   if num_bars == 1:
      plt.xlim([-0.5, 0.5])
             
   sns.set_style("whitegrid")
   sns.despine(left=True)
   plt.grid(False)
   plt.gca().spines['top'].set_visible(False)
   plt.gca().spines['right'].set_visible(False)
   plt.gca().spines['bottom'].set_color('#dddddd')
   plt.gca().spines['left'].set_color('#dddddd')
   plt.gca().xaxis.label.set_color('#333333')
   plt.gca().yaxis.label.set_color('#333333')
   plt.gca().tick_params(colors='#333333')

   img_path = f'{category}.png'
   plt.savefig(img_path, dpi=300, bbox_inches='tight')

   # Add the image to the slide
   slide = prs.slides[slide_index]
   left = Inches(0.9)
   top = Inches(1.8)
   width = Inches(5.8)
   height = Inches(3.5)
   slide.shapes.add_picture(img_path, left, top, width, height)

##------------------>> recommendations automated --------------------##
reco_ = pd.read_csv('recommendations/recommendations_Summary.csv', sep=',')
data = pd.read_csv(f'{dir_}Summary_{company}_.csv', sep=',')
data['Strategy'] = data['Strategy'].round(0)
data['Measurement'] = data['Measurement'].round(0)
data['Data'] = data['Data'].round(0)
data['Processes'] = data['Processes'].round(0)
data['Talent'] = data['Talent'].round(0)
data['Reporting'] = data['Reporting'].round(0)
data['Technology'] = data['Technology'].round(0)

categories = ['Strategy', 'Measurement', 'Data', 'Processes', 'Talent', 'Reporting', 'Technology']
recommendations = {}

for category in categories:
    if data[category].iloc[0] == 0:
        recommendations[category] = reco_[category][data[category][0]]
    else:
        recommendations[category] = reco_[category][data[category][0] - 1]

reco_ = pd.read_csv('recommendations/recommendations_Strategy.csv', sep=',')
data = pd.read_csv(f'{dir_}Strategy_{company}_.csv', sep=',')
values = data['max'][:-1]
reco_st_1 = reco_['Capability Purpose'][values[0]-1]
reco_st_2 = reco_['Capability Stakeholders'][values[1]-1]
reco_st_3 = reco_['Impact Strategy'][values[2]-1]

reco_ = pd.read_csv('recommendations/recommendations_Talent.csv', sep=',')
data = pd.read_csv(f'{dir_}Talent_{company}_.csv', sep=',')
values = data['max'][:-1]
reco_ta_1 = reco_['Equipping'][values[0]-1]
reco_ta_2 = reco_['Impact Performance'][values[1]-1]
reco_ta_3 = reco_['Team Composition'][values[2]-1]

reco_ = pd.read_csv('recommendations/recommendations_Processes.csv', sep=',')
data = pd.read_csv(f'{dir_}Processes_{company}_.csv', sep=',')
values = data['max'][:-1]
reco_pr_1 = reco_['Processes & Responsibility Framework'][values[0]-1]

reco_ = pd.read_csv('recommendations/recommendations_Data.csv', sep=',')
data = pd.read_csv(f'{dir_}Data_{company}_.csv', sep=',')
values = data['max'][:-1]
reco_da_1 = reco_['Data Access'][values[0]-1]
reco_da_2 = reco_['Data Collection'][values[1]-1]
reco_da_3 = reco_['Data Quality'][values[2]-1]

reco_ = pd.read_csv('recommendations/recommendations_Measurement.csv', sep=',')
data = pd.read_csv(f'{dir_}Measurement_{company}_.csv', sep=',')
values = data['max'][:-1]
reco_me_1 = reco_['IMF'][values[0]-1]
reco_me_2 = reco_['Tools & Templates'][values[1]-1]
reco_me_3 = reco_['Evaluation'][values[2]-1]
reco_me_4 = reco_['Research, Knowledge, & Insights'][values[3]-1]

reco_ = pd.read_csv('recommendations/recommendations_Reporting.csv', sep=',')
data = pd.read_csv(f'{dir_}Reporting_{company}_.csv', sep=',')
values = data['max'][:-1]
reco_re_1 = reco_['Reporting Framework'][values[0]-1]
reco_re_2 = reco_['Reporting Standards'][values[1]-1]

reco_ = pd.read_csv('recommendations/recommendations_Technology.csv', sep=',')
data = pd.read_csv(f'{dir_}Technology_{company}_.csv', sep=',')
values = data['max'][:-1]
reco_te_1 = reco_['Technology'][values[0]-1]

##------------------- recommendations automated <<--------------------##

   # Load the presentation
prs = Presentation('Capability Assessment Report Template-pdf.pptx')  
#company='Ehlke-Hepworth'
date_ = 'April 2024'
relativ = 'Relativ Impact'

from pptx.util import Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

##------------------->> recommendations automated --------------------##
def add_formatted_textbox1(slide, text, left, top, width, height):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.text = text
    tf.margin_top = Inches(0)
    tf.margin_bottom = Inches(0)
    tf.margin_left = Inches(0)
    tf.margin_right = Inches(0)
    tf.auto_size = None
          
    for paragraph in tf.paragraphs:
        paragraph.font.bold = False
        paragraph.font.size = Pt(9.5) #10
        paragraph.font.color.rgb = RGBColor(66, 83, 105)
        paragraph.alignment = PP_ALIGN.LEFT
        paragraph.font.name = 'Avenir Next'

# Usage
slide = prs.slides[8]
add_formatted_textbox1(slide, recommendations['Strategy'], 1.3, 1.7, 7, 1)
add_formatted_textbox1(slide, recommendations['Measurement'], 1.3, 2.2, 7, 1)
add_formatted_textbox1(slide, recommendations['Data'], 1.3, 2.7, 7, 1)
add_formatted_textbox1(slide, recommendations['Talent'], 1.3, 3.2, 7, 1)
add_formatted_textbox1(slide, recommendations['Processes'], 1.3, 3.7, 7, 1)
add_formatted_textbox1(slide, recommendations['Reporting'], 1.3, 4.2, 7, 1)
add_formatted_textbox1(slide, recommendations['Technology'], 1.3, 4.7, 7, 1)

def add_formatted_textbox(slide, text, left, top, width, height):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.text = text
    tf.text = text
    tf.margin_top = Inches(0)
    tf.margin_bottom = Inches(0)
    tf.margin_left = Inches(0)
    tf.margin_right = Inches(0)
    tf.auto_size = None
          
    for paragraph in tf.paragraphs:
        paragraph.font.bold = False
        paragraph.font.size = Pt(9.5) #10
        paragraph.font.color.rgb = RGBColor(66, 83, 105)
        paragraph.alignment = PP_ALIGN.LEFT
        paragraph.font.name = 'Avenir Next'
              
slide = prs.slides[13]
add_formatted_textbox(slide, reco_st_1, 2, 1.3, 7, 1)
add_formatted_textbox(slide, reco_st_2, 2, 2.2, 7, 1)
add_formatted_textbox(slide, reco_st_3, 2, 2.9, 7, 1)
slide = prs.slides[17]
add_formatted_textbox(slide, reco_ta_1, 2, 1.2, 7, 1)
add_formatted_textbox(slide, reco_ta_2, 2, 3.1, 7, 1)
add_formatted_textbox(slide, reco_ta_3, 2, 4.05, 7, 1)
slide = prs.slides[21]
add_formatted_textbox(slide, reco_pr_1, 2, 1.3, 7, 1)
slide = prs.slides[25]
add_formatted_textbox(slide, reco_da_1, 2, 1.3, 7, 1)
add_formatted_textbox(slide, reco_da_2, 2, 2.8, 7, 1)
add_formatted_textbox(slide, reco_da_3, 2, 4.3, 7, 1)
slide = prs.slides[29]
add_formatted_textbox(slide, reco_me_1, 2, 1.3, 7, 1)
add_formatted_textbox(slide, reco_me_2, 2, 2.9, 7, 1)
slide = prs.slides[30]
add_formatted_textbox(slide, reco_me_3, 2, 1.3, 7, 1)
add_formatted_textbox(slide, reco_me_4, 2, 2.75, 7, 1)
slide = prs.slides[34]
add_formatted_textbox(slide, reco_re_1, 2, 1.3, 7, 1)
add_formatted_textbox(slide, reco_re_2, 2, 2.1, 7, 1)
slide = prs.slides[38]
add_formatted_textbox(slide, reco_te_1, 2, 1.3, 7, 1)

##------------------- recommendations automated <<--------------------##

slide = prs.slides[0]
txBox = slide.shapes.add_textbox(Inches(6.8), Inches(0.5), Inches(4), Inches(0.5)) # Inches(8.95),
tf = txBox.text_frame
tf.text = str(company)
tf.text = str(company).upper()

tf.margin_top = Inches(0)
tf.margin_bottom = Inches(0)
tf.margin_left = Inches(0)
tf.margin_right = Inches(0)
tf.auto_size = None

p = tf.paragraphs[0]
p.font.bold = False
p.font.size = Pt(32) #32
p.font.color.rgb = RGBColor(66, 83, 105)
p.alignment = PP_ALIGN.RIGHT
p.font.name = 'Avenir Next'

# Add Relativ to the slide
slide = prs.slides[0]
txBox_date = slide.shapes.add_textbox(Inches(8.4), Inches(4), Inches(2), Inches(0.5))
tf_rel = txBox_date.text_frame
tf_rel.text = relativ
tf_rel.text = relativ.upper()

p_rel = tf_rel.paragraphs[0]
p_rel.font.bold = True
p_rel.font.size = Pt(21)
p_rel.font.color.rgb = RGBColor(66, 83, 105)
p_rel.alignment = PP_ALIGN.LEFT
p_rel.font.name = 'Avenir Next'

# Add date to the slide
slide = prs.slides[0]
txBox_date = slide.shapes.add_textbox(Inches(8.95), Inches(4.4), Inches(2), Inches(0.5))
tf_date = txBox_date.text_frame
tf_date.text = date_
tf_date.text = date_.upper()

p_date = tf_date.paragraphs[0]
p_date.font.bold = False
p_date.font.size = Pt(18)
p_date.font.color.rgb = RGBColor(66, 83, 105)
p_date.alignment = PP_ALIGN.RIGHT
p_date.font.name = 'Avenir Next'

# Add recommendations to slide 8
cleaned_words = [word.replace('\n', ' ') for word in related_words]  # Replace newline characters with spaces

slide = prs.slides[7]
txBox_rec8 = slide.shapes.add_textbox(Inches(4.4), Inches(1.16), Inches(7), Inches(0.5)) #(Inches(1.36), Inches(1.165)
tf_rec8 = txBox_rec8.text_frame
if 1 < len(cleaned_words) < 6:
   tf_rec8.text = ', '.join(cleaned_words[:-1]) + ', and ' + cleaned_words[-1] + '.'
if len(cleaned_words) > 5:
   tf_rec8.text = ', '.join(cleaned_words[:5]) +','
if len(cleaned_words) == 1:
   tf_rec8.text = cleaned_words[0] +'.'

tf_rec8.margin_top = Inches(0)
tf_rec8.margin_bottom = Inches(0)
tf_rec8.margin_left = Inches(0)
tf_rec8.margin_right = Inches(0)
tf_rec8.auto_size = None

p_rec8 = tf_rec8.paragraphs[0]
p_rec8.font.bold = False
p_rec8.font.size = Pt(9.5) #10.5
p_rec8.font.color.rgb = RGBColor(66, 83, 105)
p_rec8.alignment = PP_ALIGN.LEFT
p_rec8.font.name = 'Avenir Next'

slide = prs.slides[7]
txBox_rec82 = slide.shapes.add_textbox(Inches(0.7), Inches(1.33), Inches(7), Inches(0.5)) #0.6, 1.35
tf_rec82 = txBox_rec82.text_frame

if len(cleaned_words[5:]) > 1:
   tf_rec82.text = ', '.join(cleaned_words[5:-1]) + ', and ' + cleaned_words[-1] +'.'


tf_rec82.margin_top = Inches(0)
tf_rec82.margin_bottom = Inches(0)
tf_rec82.margin_left = Inches(0)
tf_rec82.margin_right = Inches(0)
tf_rec82.auto_size = None

p_rec82 = tf_rec82.paragraphs[0]
p_rec82.font.bold = False
p_rec82.font.size = Pt(9.5) #10.5
p_rec82.font.color.rgb = RGBColor(66, 83, 105)
p_rec82.alignment = PP_ALIGN.LEFT
p_rec82.font.name = 'Avenir Next'


if len(cleaned_words[5:]) == 1:
    tf_rec82.text = 'and ' + cleaned_words[-1] +'.'

tf_rec82.margin_top = Inches(0)
tf_rec82.margin_bottom = Inches(0)
tf_rec82.margin_left = Inches(0)
tf_rec82.margin_right = Inches(0)
tf_rec82.auto_size = None

p_rec82 = tf_rec82.paragraphs[0]
p_rec82.font.bold = False
p_rec82.font.size = Pt(9.5) #10.5
p_rec82.font.color.rgb = RGBColor(66, 83, 105)
p_rec82.alignment = PP_ALIGN.LEFT
p_rec82.font.name = 'Avenir Next'


# Set the font to be super light for all text in the slides
for slide in prs.slides:
   for shape in slide.shapes:
      if not shape.has_text_frame:
            continue
      for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
               run.font.weight = 'light'

   # Call the function for each category
dir_ = 'results/'
      #company='Ehlke-Hepworth'
categories = ['Strategy', 'Talent', 'Processes', 'Data', 'Measurement', 'Reporting', 'Technology']
slide_indices = [12, 16, 20, 24, 28, 33, 37]  # Add more slide indices as needed
for category, slide_index in zip(categories, slide_indices):
   create_plot_and_add_to_ppt(category, company, dir_, prs, slide_index)

      # Save the presentation
slide= prs.slides[6]
left = Inches(1.1)
top = Inches(2.1)
width = Inches(8.6)
height = Inches(3.5)

slide.shapes.add_picture(str(dir_)+str(company)+'heptagon_plot.png', left, top, width, height)

slide= prs.slides[7]
left = Inches(1.6)
top = Inches(1.8)
width = Inches(8.2)
height = Inches(3.65)

slide.shapes.add_picture(str(dir_)+'CapAss.png', left, top, width, height)

      
prs.save('Capability Assessment Report.pptx')
        
"""
import smtplib
from email.mime.multipart import MIMEMultipart
from email.mime.base import MIMEBase
from email.mime.application import MIMEApplication
from email.message import EmailMessage
from email import encoders

sender_email = "ehlke.hepworth@relativ.co.za"
receiver_email = "ehlke.hepworth@outlook.com"
password = os.getenv('REPORT_SECRET')
if not password:
    raise ValueError("Password not found in environment variables")


msg = MIMEMultipart()
msg['From'] = sender_email
msg['To'] = receiver_email
msg['Subject'] = "Capability Assessment Report"
#body = "Dear "+company+", \
#      \n Please find attached your Capability Assessment Report."
#msg.attach(MIMEText(body, 'plain'))

msg.attach(MIMEApplication(_data="Dear"+company+ ", \
                           \nPlease find attached your capability assessment report.", _subtype="plain"))

file_path = 'Capability Assessment Report - '+company+'_.pptx'

with open(file_path, "rb") as f:
    file_attachment = MIMEBase('application', 'octet-stream')
    file_attachment.set_payload(f.read())
    encoders.encode_base64(file_attachment)
    file_attachment.add_header('Content-Disposition', 'attachment; filename="{}"'.format(file_path.split("/")[-1]))
    msg.attach(file_attachment)

try:
    # Connect to the server using SMTP and upgrade to SSL/TLS using STARTTLS
    with smtplib.SMTP("smtp.office365.com", 587) as server:
        server.ehlo()  # Can be omitted
        server.starttls()  # Secure the connection
        server.ehlo()  # Can be omitted
        server.login(sender_email, password)
        server.send_message(message)
        print("Email sent successfully!")
except Exception as e:
    print(f"Failed to send email: {e}")
   #capability_report(company=company)
"""

