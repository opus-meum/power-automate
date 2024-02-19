##Automated recommendations in response to survey answers.
"""
Created on Mon Dec 04 2023

@authori ehlke_hepworth
"""

#%% STRATEGY

import pandas as pd
import matplotlib.pyplot as plt
import numpy as np
import csv
import pandas as pd
import glob
import os

#%%
dir_ = '/Users/ehlke/Desktop/Capability_Assessment_Results/'
data = pd.read_csv(dir_+'Capability Assessment Survey Responses_12122023.csv')

word_to_number = {'nan':0, 'Nascent': 1, 'Emerging': 2, 'Expanding': 3, 'Optimising':4, 'Mature': 5}  # Add more words and their corresponding numbers


####STRATEGY
strategy_capp = pd.read_csv(dir_+'Data/CapA_00_capabilitypurpose.csv')
strategy_caps = pd.read_csv(dir_+'Data/CapA_00_capabilitystakeholders.csv')
strategy_imps = pd.read_csv(dir_+'Data/CapA_00_imapctstrategy.csv')

sentences_list = []
for i in range(len(data)):
    row_data = data.iloc[i, 8:].apply(lambda x: str(x).split(';')).sum()
    sentences_list.append(row_data)


one=[]
two=[]
three=[]
four=[]
five=[]


for i in range(len(sentences_list)):
    oo = strategy_capp[strategy_capp.isin(sentences_list[i])].groupby(['nascent']).size().sum()
    one.append(oo)
    tw = strategy_capp[strategy_capp.isin(sentences_list[i])].groupby(['emerging']).size().sum()
    two.append(tw)
    th = strategy_capp[strategy_capp.isin(sentences_list[i])].groupby(['expanding']).size().sum()
    three.append(th)
    fo = strategy_capp[strategy_capp.isin(sentences_list[i])].groupby(['optimising']).size().sum()
    four.append(fo)
    fi = strategy_capp[strategy_capp.isin(sentences_list[i])].groupby(['mature']).size().sum()
    five.append(fi)

# Repeating the script for strategy_caps
one_caps=[]
two_caps=[]
three_caps=[]
four_caps=[]
five_caps=[]

for i in range(len(sentences_list)):
    oo_caps = strategy_caps[strategy_caps.isin(sentences_list[i])].groupby(['nascent']).size().sum()
    one_caps.append(oo_caps)
    tw_caps = strategy_caps[strategy_caps.isin(sentences_list[i])].groupby(['emerging']).size().sum()
    two_caps.append(tw_caps)
    th_caps = strategy_caps[strategy_caps.isin(sentences_list[i])].groupby(['expanding']).size().sum()
    three_caps.append(th_caps)
    fo_caps = strategy_caps[strategy_caps.isin(sentences_list[i])].groupby(['optimising']).size().sum()
    four_caps.append(fo_caps)
    fi_caps = strategy_caps[strategy_caps.isin(sentences_list[i])].groupby(['mature']).size().sum()
    five_caps.append(fi_caps)

# Repeating the script for strategy_imps
one_imps=[]
two_imps=[]
three_imps=[]
four_imps=[]
five_imps=[]

for i in range(len(sentences_list)):
    oo_imps = strategy_imps[strategy_imps.isin(sentences_list[i])].groupby(['nascent']).size().sum()
    one_imps.append(oo_imps)
    tw_imps = strategy_imps[strategy_imps.isin(sentences_list[i])].groupby(['emerging']).size().sum()
    two_imps.append(tw_imps)
    th_imps = strategy_imps[strategy_imps.isin(sentences_list[i])].groupby(['expanding']).size().sum()
    three_imps.append(th_imps)
    fo_imps = strategy_imps[strategy_imps.isin(sentences_list[i])].groupby(['optimising']).size().sum()
    four_imps.append(fo_imps)
    fi_imps = strategy_imps[strategy_imps.isin(sentences_list[i])].groupby(['mature']).size().sum()
    five_imps.append(fi_imps)


#combined_data = pd.concat([one, one_caps, one_imps], axis=1)
data_dict = {'dimension':['Capability Purpose', 'Capability Stakeholders', 'Impact Strategy'],'Nascent': [one,one_caps,one_imps], 'Emerging': [two,two_caps,two_imps], 'Expanding': [three,three_caps,three_imps],\
              'Optimising': [four,four_caps,four_imps], 'Mature': [five,five_caps,five_imps]}

df = pd.DataFrame(data_dict)

tst=[]
fn=[]
save_=[]

#for k in range(len(strategy_capp)):

for i in range(len(df)):
    tmp = df.iloc[:,1:].apply(lambda x: x.str[i])
    tst.append(tmp)
    filename = data.iloc[i, 6] + '_.csv'
    fn.append(filename)
    total_row = tst[i].sum()
    tmp = tst[i].append(total_row, ignore_index=True)    
    tmp.insert(0,'Dimension',['Capability Purpose', 'Capability Stakeholders', 'Impact Strategy','Total'])
    
    columns_greater_than_zero = tmp.iloc[:,1:].apply(lambda row: row.index[row.astype(float) > 0].tolist(), axis=1)
    columns_greater_than_zero_df = columns_greater_than_zero.apply(pd.Series)
    columns_greater_than_zero_df = columns_greater_than_zero_df.replace(word_to_number)
    average_values = columns_greater_than_zero_df.mean(axis=1)
    average_values = average_values.fillna(0)
    print(average_values)   
    tmp['level'] = average_values
    last_value_level = tmp['level'].iloc[-1]
    tmp['level'].iloc[-1] = tmp['level'].mean()
    max_columns = tmp.iloc[:, 1:-1].apply(lambda x: x.idxmax() if x.max() > 0 else 'nan', axis=1)
   # max_columns = tmp.iloc[:, 1:-1].idxmax(axis=1) 
    tmp['max_column'] = max_columns  
    print(tmp)
    tmp['max'] = tmp['max_column'].map(word_to_number)
    tmp['max'] = tmp['level'].apply(np.floor)
    tmp['idx'] = tmp.index +1

    print(tmp)

    tmp.to_csv(dir_+'Strategy_' +fn[i], index=False)
    tmp[['level'][-1]].to_frame().rename(columns={"level":'Strategy'}).tail(1).to_csv(dir_+'Summary_'+fn[i], index=False)


#-----------------------------------------------------------------------------------------------
#-----------------------------------------------------------------------------------------------
#-----------------------------------------------------------------------------------------------
####TALENT
file1 = pd.read_csv(dir_+'Data/CapA_01_equipping.csv')
file2 = pd.read_csv(dir_+'Data/CapA_01_impactperformance.csv')
file3 = pd.read_csv(dir_+'Data/CapA_01_teamcomposition.csv')


one=[]
two=[]
three=[]
four=[]
five=[]


for i in range(len(sentences_list)):
    oo = file1[file1.isin(sentences_list[i])].groupby(['nascent']).size().sum()
    one.append(oo)
    tw = file1[file1.isin(sentences_list[i])].groupby(['emerging']).size().sum()
    two.append(tw)
    th = file1[file1.isin(sentences_list[i])].groupby(['expanding']).size().sum()
    three.append(th)
    fo = file1[file1.isin(sentences_list[i])].groupby(['optimising']).size().sum()
    four.append(fo)
    fi = file1[file1.isin(sentences_list[i])].groupby(['mature']).size().sum()
    five.append(fi)

# Repeating the script for strategy_caps
one2=[]
two2=[]
three2=[]
four2=[]
five2=[]

for i in range(len(sentences_list)):
    oo_caps = file2[file2.isin(sentences_list[i])].groupby(['nascent']).size().sum()
    one2.append(oo_caps)
    tw_caps = file2[file2.isin(sentences_list[i])].groupby(['emerging']).size().sum()
    two2.append(tw_caps)
    th_caps = file2[file2.isin(sentences_list[i])].groupby(['expanding']).size().sum()
    three2.append(th_caps)
    fo_caps = file2[file2.isin(sentences_list[i])].groupby(['optimising']).size().sum()
    four2.append(fo_caps)
    fi_caps = file2[file2.isin(sentences_list[i])].groupby(['mature']).size().sum()
    five2.append(fi_caps)

one3=[]
two3=[]
three3=[]
four3=[]
five3=[]

for i in range(len(sentences_list)):
    oo_imps = file3[file3.isin(sentences_list[i])].groupby(['nascent']).size().sum()
    one3.append(oo_imps)
    tw_imps = file3[file3.isin(sentences_list[i])].groupby(['emerging']).size().sum()
    two3.append(tw_imps)
    th_imps = file3[file3.isin(sentences_list[i])].groupby(['expanding']).size().sum()
    three3.append(th_imps)
    fo_imps = file3[file3.isin(sentences_list[i])].groupby(['optimising']).size().sum()
    four3.append(fo_imps)
    fi_imps = file3[file3.isin(sentences_list[i])].groupby(['mature']).size().sum()
    five3.append(fi_imps)

#combined_data = pd.concat([one, one_caps, one_imps], axis=1)
data_dict = {'dimension':['Equipping', 'Impact Performance', 'Team Composition'],'Nascent': [one,one2,one3], 'Emerging': [two,two2,two3], 'Expanding': [three,three2,three3],\
              'Optimising': [four,four2,four3], 'Mature': [five,five2,five3]}

df = pd.DataFrame(data_dict)

tst=[]
fn=[]
save_=[]
summary = []

#for k in range(len(strategy_capp)):

for i in range(len(df)):
    tmp = df.iloc[:,1:].apply(lambda x: x.str[i])
    tst.append(tmp)
    filename = data.iloc[i, 6] + '_.csv'
    fn.append(filename)
    total_row = tst[i].sum()
    tmp = tst[i].append(total_row, ignore_index=True)
    tmp.insert(0,'Dimension',['Equipping', 'Impact Performance', 'Team Composition','Total'])

    columns_greater_than_zero = tmp.iloc[:,1:].apply(lambda row: row.index[row.astype(float) > 0].tolist(), axis=1)
    columns_greater_than_zero_df = columns_greater_than_zero.apply(pd.Series)
    columns_greater_than_zero_df = columns_greater_than_zero_df.replace(word_to_number)
    average_values = columns_greater_than_zero_df.mean(axis=1)
    average_values = average_values.fillna(0)

    tmp['level'] = average_values
    last_value_level = tmp['level'].iloc[-1]
    tmp['level'].iloc[-1] = tmp['level'].mean()
        
    max_columns = tmp.iloc[:, 1:-1].apply(lambda x: x.idxmax() if x.max() > 0 else 'nan', axis=1)
    tmp['max_column'] = max_columns  
    tmp['max'] = tmp['max_column'].map(word_to_number)

    tmp['max'] = tmp['level'].apply(np.floor)
    tmp['idx'] = tmp.index +1

    tmp.to_csv(dir_+'Talent_' +fn[i], index=False)
    print(tmp)
    add_ = tmp[['level'][-1]].to_frame().rename(columns={"level":'Talent'}).tail(1)
    add_=pd.DataFrame(add_)
    print(add_['Talent'].values[0])

    summary_ = pd.read_csv(dir_+'Summary_'+fn[i])
    print(summary_)
    summary_['Talent'] = add_['Talent'].values[0]
    print(summary_)
    summary_.to_csv(dir_+'Summary_' +fn[i], index=False)



#-----------------------------------------------------------------------------------------------
#-----------------------------------------------------------------------------------------------
#-----------------------------------------------------------------------------------------------
####PROCESSES
file1 = pd.read_csv(dir_+'Data/CapA_02_processes.csv')
file2 = pd.read_csv(dir_+'Data/CapA_02_responsibilityframework.csv')


one=[]
two=[]
three=[]
four=[]
five=[]


for i in range(len(sentences_list)):
    oo = file1[file1.isin(sentences_list[i])].groupby(['nascent']).size().sum()
    one.append(oo)
    tw = file1[file1.isin(sentences_list[i])].groupby(['emerging']).size().sum()
    two.append(tw)
    th = file1[file1.isin(sentences_list[i])].groupby(['expanding']).size().sum()
    three.append(th)
    fo = file1[file1.isin(sentences_list[i])].groupby(['optimising']).size().sum()
    four.append(fo)
    fi = file1[file1.isin(sentences_list[i])].groupby(['mature']).size().sum()
    five.append(fi)

one2=[]
two2=[]
three2=[]
four2=[]
five2=[]

for i in range(len(sentences_list)):
    oo_caps = file2[file2.isin(sentences_list[i])].groupby(['nascent']).size().sum()
    one2.append(oo_caps)
    tw_caps = file2[file2.isin(sentences_list[i])].groupby(['emerging']).size().sum()
    two2.append(tw_caps)
    th_caps = file2[file2.isin(sentences_list[i])].groupby(['expanding']).size().sum()
    three2.append(th_caps)
    fo_caps = file2[file2.isin(sentences_list[i])].groupby(['optimising']).size().sum()
    four2.append(fo_caps)
    fi_caps = file2[file2.isin(sentences_list[i])].groupby(['mature']).size().sum()
    five2.append(fi_caps)


#combined_data = pd.concat([one, one_caps, one_imps], axis=1)
data_dict = {'dimension':['Processes', 'Responsibility \nFramework'],'Nascent': [one,one2], 'Emerging': [two,two2], 'Expanding': [three,three2],\
              'Optimising': [four,four2], 'Mature': [five,five2]}

df = pd.DataFrame(data_dict)

tst=[]
fn=[]
save_=[]

#for k in range(len(strategy_capp)):

for i in range(len(df)):
    tmp = df.iloc[:,1:].apply(lambda x: x.str[i])
    tst.append(tmp)
    filename = data.iloc[i, 6] + '_.csv'
    fn.append(filename)
    total_row = tst[i].sum()
    print(total_row)
    tmp = tst[i].append(total_row, ignore_index=True)
    tmp.insert(0,'Dimension',['Processes', 'Responsibility \nFramework','Total'])

    columns_greater_than_zero = tmp.iloc[:,1:].apply(lambda row: row.index[row.astype(float) > 0].tolist(), axis=1)
    columns_greater_than_zero_df = columns_greater_than_zero.apply(pd.Series)
    columns_greater_than_zero_df = columns_greater_than_zero_df.replace(word_to_number)
    average_values = columns_greater_than_zero_df.mean(axis=1)
    average_values = average_values.fillna(0)
    print(average_values)   

    tmp['level'] = average_values
    last_value_level = tmp['level'].iloc[-1]
    tmp['level'].iloc[-1] = tmp['level'].mean()

    max_columns = tmp.iloc[:, 1:-1].apply(lambda x: x.idxmax() if x.max() > 0 else 'nan', axis=1)
    tmp['max_column'] = max_columns  
    tmp['max'] = tmp['max_column'].map(word_to_number)
    tmp['max'] = tmp['level'].apply(np.floor)
    tmp['idx'] = tmp.index +1

    tmp.to_csv(dir_+'Processes_' +fn[i], index=False)
    print(tmp)
    add_ = tmp[['level'][-1]].to_frame().rename(columns={"level":'Processes'}).tail(1)
    add_=pd.DataFrame(add_)
    print(add_['Processes'].values[0])

    summary_ = pd.read_csv(dir_+'Summary_'+fn[i])
    print(summary_)
    summary_['Processes'] = add_['Processes'].values[0]
    print(summary_)
    summary_.to_csv(dir_+'Summary_' +fn[i], index=False)



#-----------------------------------------------------------------------------------------------
#-----------------------------------------------------------------------------------------------
#-----------------------------------------------------------------------------------------------
####DATA
file1 = pd.read_csv(dir_+'Data/CapA_03_dataaccess.csv')
file2 = pd.read_csv(dir_+'Data/CapA_03_datacollection.csv')
file3 = pd.read_csv(dir_+'Data/CapA_03_dataquality.csv')


one=[]
two=[]
three=[]
four=[]
five=[]


for i in range(len(sentences_list)):
    oo = file1[file1.isin(sentences_list[i])].groupby(['nascent']).size().sum()
    one.append(oo)
    tw = file1[file1.isin(sentences_list[i])].groupby(['emerging']).size().sum()
    two.append(tw)
    th = file1[file1.isin(sentences_list[i])].groupby(['expanding']).size().sum()
    three.append(th)
    fo = file1[file1.isin(sentences_list[i])].groupby(['optimising']).size().sum()
    four.append(fo)
    fi = file1[file1.isin(sentences_list[i])].groupby(['mature']).size().sum()
    five.append(fi)

one2=[]
two2=[]
three2=[]
four2=[]
five2=[]

for i in range(len(sentences_list)):
    oo_caps = file2[file2.isin(sentences_list[i])].groupby(['nascent']).size().sum()
    one2.append(oo_caps)
    tw_caps = file2[file2.isin(sentences_list[i])].groupby(['emerging']).size().sum()
    two2.append(tw_caps)
    th_caps = file2[file2.isin(sentences_list[i])].groupby(['expanding']).size().sum()
    three2.append(th_caps)
    fo_caps = file2[file2.isin(sentences_list[i])].groupby(['optimising']).size().sum()
    four2.append(fo_caps)
    fi_caps = file2[file2.isin(sentences_list[i])].groupby(['mature']).size().sum()
    five2.append(fi_caps)

one3=[]
two3=[]
three3=[]
four3=[]
five3=[]

for i in range(len(sentences_list)):
    oo_imps = file3[file3.isin(sentences_list[i])].groupby(['nascent']).size().sum()
    one3.append(oo_imps)
    tw_imps = file3[file3.isin(sentences_list[i])].groupby(['emerging']).size().sum()
    two3.append(tw_imps)
    th_imps = file3[file3.isin(sentences_list[i])].groupby(['expanding']).size().sum()
    three3.append(th_imps)
    fo_imps = file3[file3.isin(sentences_list[i])].groupby(['optimising']).size().sum()
    four3.append(fo_imps)
    fi_imps = file3[file3.isin(sentences_list[i])].groupby(['mature']).size().sum()
    five3.append(fi_imps)


#combined_data = pd.concat([one, one_caps, one_imps], axis=1)
data_dict = {'dimension':['Data Access', 'Data Collection', 'Data Quality'],'Nascent': [one,one2,one3], 'Emerging': [two,two2,two3], 'Expanding': [three,three2,three3],\
              'Optimising': [four,four2,four3], 'Mature': [five,five2,five3]}

df = pd.DataFrame(data_dict)

tst=[]
fn=[]
save_=[]

#for k in range(len(strategy_capp)):

for i in range(len(df)):
    tmp = df.iloc[:,1:].apply(lambda x: x.str[i])
    tst.append(tmp)
    filename = data.iloc[i, 6] + '_.csv'
    fn.append(filename)
    total_row = tst[i].sum()
    print(total_row)
    tmp = tst[i].append(total_row, ignore_index=True)
    tmp.insert(0,'Dimension',['Data Access', 'Data Collection', 'Data Quality','Total'])
    print(tmp)

    columns_greater_than_zero = tmp.iloc[:,1:].apply(lambda row: row.index[row.astype(float) > 0].tolist(), axis=1)
    columns_greater_than_zero_df = columns_greater_than_zero.apply(pd.Series)
    columns_greater_than_zero_df = columns_greater_than_zero_df.replace(word_to_number)
    average_values = columns_greater_than_zero_df.mean(axis=1)
    average_values = average_values.fillna(0)
    print(average_values)   
    tmp['level'] = average_values

    max_columns = tmp.iloc[:, 1:-1].apply(lambda x: x.idxmax() if x.max() > 0 else 'nan', axis=1)
    tmp['max_column'] = max_columns  
    tmp['max'] = tmp['max_column'].map(word_to_number)
    tmp['max'] = tmp['level'].apply(np.floor)
    tmp['idx'] = tmp.index +1


    tmp.to_csv(dir_+'Data_' +fn[i], index=False)

    add_ = tmp[['level'][-1]].to_frame().rename(columns={"level":'Data'}).tail(1)
    add_=pd.DataFrame(add_)
    print(add_['Data'].values[0])

    summary_ = pd.read_csv(dir_+'Summary_'+fn[i])
    print(summary_)
    summary_['Data'] = add_['Data'].values[0]
    print(summary_)
    summary_.to_csv(dir_+'Summary_' +fn[i], index=False)

#-----------------------------------------------------------------------------------------------
#-----------------------------------------------------------------------------------------------
#-----------------------------------------------------------------------------------------------
####REPORTING  
file1 = pd.read_csv(dir_+'Data/CapA_05_framework.csv')
file2 = pd.read_csv(dir_+'Data/CapA_05_standards.csv')

one=[]
two=[]
three=[]
four=[]
five=[]


for i in range(len(sentences_list)):
    oo = file1[file1.isin(sentences_list[i])].groupby(['nascent']).size().sum()
    one.append(oo)
    tw = file1[file1.isin(sentences_list[i])].groupby(['emerging']).size().sum()
    two.append(tw)
    th = file1[file1.isin(sentences_list[i])].groupby(['expanding']).size().sum()
    three.append(th)
    fo = file1[file1.isin(sentences_list[i])].groupby(['optimising']).size().sum()
    four.append(fo)
    fi = file1[file1.isin(sentences_list[i])].groupby(['mature']).size().sum()
    five.append(fi)

one2=[]
two2=[]
three2=[]
four2=[]
five2=[]

for i in range(len(sentences_list)):
    oo_caps = file2[file2.isin(sentences_list[i])].groupby(['nascent']).size().sum()
    one2.append(oo_caps)
    tw_caps = file2[file2.isin(sentences_list[i])].groupby(['emerging']).size().sum()
    two2.append(tw_caps)
    th_caps = file2[file2.isin(sentences_list[i])].groupby(['expanding']).size().sum()
    three2.append(th_caps)
    fo_caps = file2[file2.isin(sentences_list[i])].groupby(['optimising']).size().sum()
    four2.append(fo_caps)
    fi_caps = file2[file2.isin(sentences_list[i])].groupby(['mature']).size().sum()
    five2.append(fi_caps)


#combined_data = pd.concat([one, one_caps, one_imps], axis=1)
data_dict = {'dimension':['Reporting Framework', 'Reporting Standards'],'Nascent': [one,one2], 'Emerging': [two,two2], 'Expanding': [three,three2],\
              'Optimising': [four,four2], 'Mature': [five,five2]}

df = pd.DataFrame(data_dict)

tst=[]
fn=[]
save_=[]

#for k in range(len(strategy_capp)):

for i in range(len(df)):
    tmp = df.iloc[:,1:].apply(lambda x: x.str[i])
    tst.append(tmp)
    filename = data.iloc[i, 6] + '_.csv'
    fn.append(filename)
    total_row = tst[i].sum()
    print(total_row)
    tmp = tst[i].append(total_row, ignore_index=True)
    tmp.insert(0,'Dimension',['Reporting Framework', 'Reporting Standards','Total'])
    print(tmp)

    columns_greater_than_zero = tmp.iloc[:,1:].apply(lambda row: row.index[row.astype(float) > 0].tolist(), axis=1)
    columns_greater_than_zero_df = columns_greater_than_zero.apply(pd.Series)
    columns_greater_than_zero_df = columns_greater_than_zero_df.replace(word_to_number)
    average_values = columns_greater_than_zero_df.mean(axis=1)
    average_values = average_values.fillna(0)
    print(average_values)   

    tmp['level'] = average_values
    last_value_level = tmp['level'].iloc[-1]
    tmp['level'].iloc[-1] = tmp['level'].mean()

    max_columns = tmp.iloc[:, 1:-1].apply(lambda x: x.idxmax() if x.max() > 0 else 'nan', axis=1)
    tmp['max_column'] = max_columns  
    tmp['max'] = tmp['max_column'].map(word_to_number)
    tmp['max'] = tmp['level'].apply(np.floor)
    tmp['idx'] = tmp.index +1

    tmp.to_csv(dir_+'Reporting_' +fn[i], index=False)

    print(tmp)
    add_ = tmp[['level'][-1]].to_frame().rename(columns={"level":'Reporting'}).tail(1)
    add_=pd.DataFrame(add_)
    print(add_['Reporting'].values[0])

    summary_ = pd.read_csv(dir_+'Summary_'+fn[i])
    print(summary_)
    summary_['Reporting'] = add_['Reporting'].values[0]
    print(summary_)
    summary_.to_csv(dir_+'Summary_' +fn[i], index=False)


#-----------------------------------------------------------------------------------------------
#-----------------------------------------------------------------------------------------------
#-----------------------------------------------------------------------------------------------
####TECHNOLOGY  
file1 = pd.read_csv(dir_+'Data/CapA_06_technology.csv')


one=[]
two=[]
three=[]
four=[]
five=[]


for i in range(len(sentences_list)):
    oo = file1[file1.isin(sentences_list[i])].groupby(['nascent']).size().sum()
    one.append(oo)
    tw = file1[file1.isin(sentences_list[i])].groupby(['emerging']).size().sum()
    two.append(tw)
    th = file1[file1.isin(sentences_list[i])].groupby(['expanding']).size().sum()
    three.append(th)
    fo = file1[file1.isin(sentences_list[i])].groupby(['optimising']).size().sum()
    four.append(fo)
    fi = file1[file1.isin(sentences_list[i])].groupby(['mature']).size().sum()
    five.append(fi)



#combined_data = pd.concat([one, one_caps, one_imps], axis=1)
data_dict = {'dimension':['Techology'],'Nascent': [one], 'Emerging': [two], 'Expanding': [three],\
              'Optimising': [four], 'Mature': [five]}

df = pd.DataFrame(data_dict)
print(len(data))

tst=[]
fn=[]
save_=[]

#for k in range(len(strategy_capp)):

for i in range(len(data)):
    tmp = df.iloc[:,1:].apply(lambda x: x.str[i])
    tst.append(tmp)
    filename = data.iloc[i, 6] + '_.csv'
    fn.append(filename)
    total_row = tst[i].sum()
    print(total_row)
    tmp = tst[i].append(total_row, ignore_index=True)
    tmp.insert(0,'Dimension',['Techology','Total'])

    columns_greater_than_zero = tmp.iloc[:,1:].apply(lambda row: row.index[row.astype(float) > 0].tolist(), axis=1)
    columns_greater_than_zero_df = columns_greater_than_zero.apply(pd.Series)
    columns_greater_than_zero_df = columns_greater_than_zero_df.replace(word_to_number)
    average_values = columns_greater_than_zero_df.mean(axis=1)
    average_values = average_values.fillna(0)
    print(average_values)   

    tmp['level'] = average_values
    last_value_level = tmp['level'].iloc[-1]
    tmp['level'].iloc[-1] = tmp['level'].mean()

    max_columns = tmp.iloc[:, 1:-1].apply(lambda x: x.idxmax() if x.max() > 0 else 'nan', axis=1)
    tmp['max_column'] = max_columns  
    tmp['max'] = tmp['max_column'].map(word_to_number)
    tmp['max'] = tmp['level'].apply(np.floor)
    tmp['idx'] = tmp.index +1

    tmp.to_csv(dir_+'Technology_' +fn[i], index=False)

    print(tmp)
    
    add_ = tmp[['level'][-1]].to_frame().rename(columns={"level":'Technology'}).tail(1)
    add_=pd.DataFrame(add_)
    print(add_['Technology'].values[0])
    
    summary_ = pd.read_csv(dir_+'Summary_'+fn[i])
    print(summary_)
    summary_['Technology'] = add_['Technology'].values[0]
    print(summary_)
    summary_.to_csv(dir_+'Summary_' +fn[i], index=False)


#-----------------------------------------------------------------------------------------------
#-----------------------------------------------------------------------------------------------
#-----------------------------------------------------------------------------------------------
####MEASUREMENT

file1 = pd.read_csv(dir_+'Data/CapA_04_IMF.csv')
file2 = pd.read_csv(dir_+'Data/CapA_04_TT.csv')
file3 = pd.read_csv(dir_+'Data/CapA_04_evaluation.csv')
file4 = pd.read_csv(dir_+'Data/CapA_04_rki.csv')


one=[]
two=[]
three=[]
four=[]
five=[]


for i in range(len(sentences_list)):
    oo = file1[file1.isin(sentences_list[i])].groupby(['nascent']).size().sum()
    one.append(oo)
    tw = file1[file1.isin(sentences_list[i])].groupby(['emerging']).size().sum()
    two.append(tw)
    th = file1[file1.isin(sentences_list[i])].groupby(['expanding']).size().sum()
    three.append(th)
    fo = file1[file1.isin(sentences_list[i])].groupby(['optimising']).size().sum()
    four.append(fo)
    fi = file1[file1.isin(sentences_list[i])].groupby(['mature']).size().sum()
    five.append(fi)

one2=[]
two2=[]
three2=[]
four2=[]
five2=[]

for i in range(len(sentences_list)):
    oo_caps = file2[file2.isin(sentences_list[i])].groupby(['nascent']).size().sum()
    one2.append(oo_caps)
    tw_caps = file2[file2.isin(sentences_list[i])].groupby(['emerging']).size().sum()
    two2.append(tw_caps)
    th_caps = file2[file2.isin(sentences_list[i])].groupby(['expanding']).size().sum()
    three2.append(th_caps)
    fo_caps = file2[file2.isin(sentences_list[i])].groupby(['optimising']).size().sum()
    four2.append(fo_caps)
    fi_caps = file2[file2.isin(sentences_list[i])].groupby(['mature']).size().sum()
    five2.append(fi_caps)

one3=[]
two3=[]
three3=[]
four3=[]
five3=[]

for i in range(len(sentences_list)):
    oo_imps = file3[file3.isin(sentences_list[i])].groupby(['nascent']).size().sum()
    one3.append(oo_imps)
    tw_imps = file3[file3.isin(sentences_list[i])].groupby(['emerging']).size().sum()
    two3.append(tw_imps)
    th_imps = file3[file3.isin(sentences_list[i])].groupby(['expanding']).size().sum()
    three3.append(th_imps)
    fo_imps = file3[file3.isin(sentences_list[i])].groupby(['optimising']).size().sum()
    four3.append(fo_imps)
    fi_imps = file3[file3.isin(sentences_list[i])].groupby(['mature']).size().sum()
    five3.append(fi_imps)

one4=[]
two4=[]
three4=[]
four4=[]
five4=[]

for i in range(len(sentences_list)):
    oo_imps = file4[file4.isin(sentences_list[i])].groupby(['nascent']).size().sum()
    one4.append(oo_imps)
    tw_imps = file4[file4.isin(sentences_list[i])].groupby(['emerging']).size().sum()
    two4.append(tw_imps)
    th_imps = file4[file4.isin(sentences_list[i])].groupby(['expanding']).size().sum()
    three4.append(th_imps)
    fo_imps = file4[file4.isin(sentences_list[i])].groupby(['optimising']).size().sum()
    four4.append(fo_imps)
    fi_imps = file4[file4.isin(sentences_list[i])].groupby(['mature']).size().sum()
    five4.append(fi_imps)


#combined_data = pd.concat([one, one_caps, one_imps], axis=1)
data_dict = {'dimension':['IMF', 'Tools & Templates', 'Evaulation', 'Risks, Knowledge,\n& Insights'],'Nascent': [one,one2,one3,one4], 'Emerging': [two,two2,two3,two4], 'Expanding': [three,three2,three3,three4],\
              'Optimising': [four,four2,four3,four4], 'Mature': [five,five,five3,five4]}

df = pd.DataFrame(data_dict)
print(df)

tst=[]
fn=[]
save_=[]

#for k in range(len(strategy_capp)):

for i in range(len(df)):
    tmp = df.iloc[:,1:].apply(lambda x: x.str[i])
    tst.append(tmp)
    filename = data.iloc[i, 6] + '_.csv'
    print(filename)
    fn.append(filename)
    print(fn)
    total_row = tst[i].sum()

    print(total_row)
    tmp = tst[i].append(total_row, ignore_index=True)
    tmp.insert(0,'Dimension',['IMF', 'Tools & Templates', 'Evaulation', 'Risks, Knowledge,\n& Insights','Total'])

    columns_greater_than_zero = tmp.iloc[:,1:].apply(lambda row: row.index[row.astype(float) > 0].tolist(), axis=1)
    columns_greater_than_zero_df = columns_greater_than_zero.apply(pd.Series)
    columns_greater_than_zero_df = columns_greater_than_zero_df.replace(word_to_number)
    average_values = columns_greater_than_zero_df.mean(axis=1)
    average_values = average_values.fillna(0)
    print(average_values)   

    tmp['level'] = average_values
    last_value_level = tmp['level'].iloc[-1]
    tmp['level'].iloc[-1] = tmp['level'].mean()

    max_columns = tmp.iloc[:, 1:-1].apply(lambda x: x.idxmax() if x.max() > 0 else 'nan', axis=1)
    tmp['max_column'] = max_columns  
    tmp['max'] = tmp['max_column'].map(word_to_number)
    tmp['max'] = tmp['level'].apply(np.floor)
    tmp['idx'] = tmp.index +1


    tmp.to_csv(dir_+'Measurement_' +fn[i], index=False)

    print(tmp)
    add_ = tmp[['level'][-1]].to_frame().rename(columns={"level":'Measurement'}).tail(1)
    add_=pd.DataFrame(add_)
    print(add_['Measurement'].values[0])

    summary_ = pd.read_csv(dir_+'Summary_'+fn[i])
    print(summary_)
    summary_['Measurement'] = add_['Measurement'].values[0]
    print(summary_)
    summary_.to_csv(dir_+'Summary_' +fn[i], index=False)



'''
 #   tst[i] = tst[i].iloc[-1,1:]
       save_.append(tst[i])

    ##process for pie chart
    tst[i] = tst[i].iloc[-1,1:]
    tst[i] = tst[i][tst[i] != 0]
    save_.append(tst[i])

    plt.figure()
    plt.title(data.iloc[i, 6]) 
    plt.pie(tst[i], labels=tst[i].index, autopct='%1.1f%%')
    plt.axis('equal')  # Equal aspect ratio ensures that pie is drawn as a circle.
    plt.savefig(dir_+data.iloc[i, 6]+'.png', dpi=300)

    print(list(tst[i]))
'''


"""
Created on Mon Nov 13 2023

@author: ehlke_hepworth
"""
#company = 'EatYourGreens'
def capability_report(company):

    import numpy as np
    import matplotlib.pyplot as plt
    import matplotlib.font_manager as fm
    import pandas as pd


    dir_ = '/Users/ehlke/Desktop/Capability_Assessment_Results/'
    df = pd.read_csv(dir_+'Summary_'+company+'_.csv')

    fig1, ax = plt.subplots(1, figsize =(17, 10))

##set up heptagon
    charcoal_color = '#425369'


    a1 = 1.2
    a2 = 0.9
    a3 = 0.6
    a4 = 0.3

    x1 = [a1 * np.cos(i * 2 * np.pi / 7) for i in range(8)]
    y1 = [a1 * np.sin(i * 2 * np.pi / 7) for i in range(8)]

    x2 = [a2 * np.cos(i * 2 * np.pi / 7) for i in range(8)]
    y2 = [a2 * np.sin(i * 2 * np.pi / 7) for i in range(8)]

    x3 = [a3 * np.cos(i * 2 * np.pi / 7) for i in range(8)]
    y3 = [a3 * np.sin(i * 2 * np.pi / 7) for i in range(8)]

    x4 = [a4 * np.cos(i * 2 * np.pi / 7) for i in range(8)]
    y4 = [a4 * np.sin(i * 2 * np.pi / 7) for i in range(8)]

    x_w = [1.77, 1, -0.3337814009344715, -1.65, -1.6,\
          -0.3337814009344719, 1,1]
    y_w = [0, 1.25, 1.53, 0.6508256086763373, -0.650825608676337,\
          -1.53, -1.23,0]

    words = ['PROCESSES', 'TALENT','STRATEGY','TECHNOLOGY','REPORTING','MEASUREMENT','DATA']

    corners = []
    for i in range(8):
        x = [a * np.cos(i * 2 * np.pi / 7) for a in [a1, a2, a3, a4]]
        y = [a * np.sin(i * 2 * np.pi / 7) for a in [a1, a2, a3, a4]]
        corners.append(list(zip(x, y)))
        x_coordinates_1 = [1.2, 0.7481877622304803, -0.26702512074757717, -1.0811626414829028, -1.081162641482903, -0.2670251207475775, 0.74818776223048, 1.2]
        y_coordinates_1 = [0.0, 0.9381977789616357, 1.1699134946181884, 0.5206604869410698, -0.5206604869410696, -1.1699134946181884, -0.9381977789616358, -2.9391523179536476e-16]

    if df['Processes'].iloc[0] == 4:
        x_pr = x1[0];
        y_pr = y1[0]
    if 3< df['Processes'].iloc[0] < 4:
        x_pr = (x1[0]+x2[0])/2;
        y_pr = (y1[0]+y2[0])/2
    if df['Processes'].iloc[0] == 3:
        x_pr = x2[0];
        y_pr = y2[0]
    if 2 < df['Processes'].iloc[0] <3:
        x_pr = (x2[0]+x3[0])/2;
        y_pr = (y2[0]+y3[0])/2
    if df['Processes'].iloc[0] == 2:
        x_pr = x3[0];
        y_pr = y3[0]
    if 1 < df['Processes'].iloc[0] <2 :
        x_pr = (x3[0]+x4[0])/2;
        y_pr = (y3[0]+y4[0])/2
    if df['Processes'].iloc[0] ==1:
        x_pr = x4[0];
        y_pr = y4[0]
    if 0 < df['Processes'].iloc[0] <1 :
        x_pr = (x4[0])/2;
        y_pr = (y4[0])/2
    if df['Processes'].iloc[0] == 0:
       x_pr = 0;
       y_pr = 0
   

    if df['Talent'].iloc[0] == 4:
      x_ta = x1[1];
      y_ta = y1[1]
    if 3 < df['Talent'].iloc[0] < 4:
      x_ta = (x1[1]+x2[1])/2;
      y_ta = (y1[1]+y2[1])/2
    if df['Talent'].iloc[0] == 3:
      x_ta = x2[1];
      y_ta = y2[1]
    if 2 < df['Talent'].iloc[0] < 3:
      x_ta = (x2[1]+x3[1])/2;
      y_ta = (y2[1]+y3[1])/2
    if df['Talent'].iloc[0] == 2:
      x_ta = x3[1];
      y_ta = y3[1]
    if 1 < df['Talent'].iloc[0] <2 :
      x_ta = (x3[1]+x4[1])/2;
      y_ta = (y3[1]+y4[1])/2
    if df['Talent'].iloc[0] ==1:
      x_ta = x4[1];
      y_ta = y4[1]
    if 0 < df['Talent'].iloc[0] <1 :
      x_ta = (x4[1])/2;
      y_ta = (y4[1])/2
    if df['Talent'].iloc[0] == 0:
      x_ta = 0;
      y_ta = 0


      
    if df['Strategy'].iloc[0] == 4:
       x_st = x1[2];
       y_st = y1[2]
    if 3 < df['Strategy'].iloc[0] < 4:
       x_st = (x1[2]+x2[2])/2;
       y_st = (y1[2]+y2[2])/2
    if df['Strategy'].iloc[0] == 3:
       x_st = x2[2];
       y_st = y2[2]
    if 2 < df['Strategy'].iloc[0] < 3:
       x_st = (x2[2]+x3[2])/2;
       y_st = (y2[2]+y3[2])/2
    if df['Strategy'].iloc[0] == 2:
       x_st = x3[2];
       y_st = y3[2]
    if 1 < df['Strategy'].iloc[0] <2 :
       x_st = (x3[2]+x4[2])/2;
       y_st = (y3[2]+y4[2])/2
    if df['Strategy'].iloc[0] ==1:
       x_st = x4[2];
       y_st = y4[2]
    if 0 < df['Strategy'].iloc[0] <1 :
       x_st = (x4[2])/2;
       y_st = (y4[2])/2
    if df['Strategy'].iloc[0] == 0:
       x_st = 0;
       y_st = 0
 
    if df['Technology'].iloc[0] == 4:
       x_te = x1[3];
       y_te = y1[3]
    if df['Technology'].iloc[0] > 3 and df['Technology'].iloc[0] < 4:
       x_te = (x1[3]+x2[3])/2;
       y_te = (y1[3]+y2[3])/2
    if df['Technology'].iloc[0] == 3:
       x_te = x2[3];
       y_te = y2[3]
    if 2 < df['Technology'].iloc[0] < 3:
       x_te = (x2[3]+x3[3])/2;
       y_te = (y2[3]+y3[3])/2
    if df['Technology'].iloc[0] == 2:
       x_te = x3[3];
       y_te = y3[3]
    if 1 < df['Technology'].iloc[0] <2 :
       x_te = (x3[3]+x4[3])/2;
       y_te = (y3[3]+y4[3])/2
    if df['Technology'].iloc[0] ==1:
       x_te = x4[3];
       y_te = y4[3]
    if 0 < df['Technology'].iloc[0] <1 :
       x_te = (x4[3])/2;
       y_te = (y4[3])/2
    if df['Technology'].iloc[0] == 0:
       x_te = 0;
       y_te = 0

    if df['Reporting'].iloc[0] == 4:
       x_re = x1[4];
       y_re = y1[4]
    if 3 < df['Reporting'].iloc[0] < 4:
       x_re = (x1[4]+x2[4])/2;
       y_re = (y1[4]+y2[4])/2
    if df['Reporting'].iloc[0] == 3:
       x_re = x2[4];
       y_re = y2[4]
    if 2 < df['Reporting'].iloc[0] < 3:
       x_re = (x2[4]+x3[4])/2;
       y_re = (y2[4]+y3[4])/2
    if df['Reporting'].iloc[0] == 2:
       x_re = x3[4];
       y_re = y3[4]
    if 1 < df['Reporting'].iloc[0] <2 :
       x_re = (x3[4]+x4[4])/2;
       y_re = (y3[4]+y4[4])/2
    if df['Reporting'].iloc[0] == 1:
       x_re = x4[4];
       y_re = y4[4]
    if 0 < df['Reporting'].iloc[0] <1 :
       x_re = (x4[4])/2;
       y_re = (y4[4])/2
    if df['Reporting'].iloc[0] == 0:
       x_re = 0;
       y_re = 0

    if df['Measurement'].iloc[0] == 4:
       x_me = x1[5];
       y_me = y1[5]
    if 3 < df['Measurement'].iloc[0] < 4:
       x_me = (x1[5]+x2[5])/2;
       y_me = (y1[5]+y2[5])/2
    if df['Measurement'].iloc[0] == 3:
       x_me = x2[5];
       y_me = y2[5]
    if 2 < df['Measurement'].iloc[0] < 3:
       x_me = (x2[5]+x3[5])/2;
       y_me = (y2[5]+y3[5])/2
    if df['Measurement'].iloc[0] == 2:
       x_me = x3[5];
       y_me = y3[5]
    if 1 < df['Measurement'].iloc[0] <2 :
       x_me = (x3[5]+x4[5])/2;
       y_me = (y3[5]+y4[5])/2
    if df['Measurement'].iloc[0] == 1:
       x_me = x4[5];
       y_me = y4[5]
    if 0 < df['Processes'].iloc[0] <1 :
       x_me = (x4[5])/2;
       y_me = (y4[5])/2
    if df['Measurement'].iloc[0] == 0:
       x_me = 0;
       y_me = 0

    if df['Data'].iloc[0] == 4:
       x_da = x1[6];
       y_da = y1[6]
    if 3 < df['Data'].iloc[0] < 4:
       x_da = (x1[6]+x2[6])/2;
       y_da = (y1[6]+y2[6])/2
    if df['Data'].iloc[0] == 3:
       x_da = x2[6];
       y_da = y2[6]
    if 2 < df['Data'].iloc[0] < 3:
       x_da = (x2[6]+x3[6])/2;
       y_da = (y2[6]+y3[6])/2
    if df['Data'].iloc[0] == 2:
       x_da = x3[6];
       y_da = y3[6]
    if 1 < df['Data'].iloc[0] <2 :
       x_da = (x3[6]+x4[6])/2;
       y_da = (y3[6]+y4[6])/2
    if df['Data'].iloc[0] == 1:
       x_da = x4[6];
       y_da = y4[6]
    if 0 < df['Processes'].iloc[0] <1 :
       x_da = (x4[6])/2;
       y_da = (y4[6])/2
    if df['Data'].iloc[0] == 0:
       x_da = 0;
       y_da = 0
# Print the coordinates for all corners of the heptagon

####set up current level of company
    x_current = [x_pr,x_ta,x_st,x_te,x_re,x_me,x_da,x_pr]
    y_current = [y_pr,y_ta,y_st,y_te,y_re,y_me,y_da,y_pr]
#x_future  = [0.901,0.568, -.226, -.535,-.654,-.160,.613,.901]
#y_future  = [0.0009,0.693,1.010,.249,-.330,-.731 ,-.772,.0009]

    x = [1.5, 0.9352347027881004, -0.3337814009344715, -1.3514533018536286, -1.3514533018536288,\
          -0.3337814009344719, 0.9352347027881001,1.5]
    y = [0, 1.1727472237020446, 1.4623918682727355, 0.6508256086763373, -0.650825608676337,\
          -1.4623918682727355, -1.1727472237020449,0]
    font_path = fm.findfont(fm.FontProperties(family='Avenir Next light'))

    # Plot the heptagon
    plt.plot(x,y, color='grey', linewidth=.5)
    plt.plot(x1,y1, color='grey',linewidth=.5)
    plt.plot(x2,y2, color='grey',linewidth=.5)
    plt.plot(x3,y3, color='grey',linewidth=.5)
    plt.plot(x4,y4, color='grey',linewidth=.5, label ="Maturity Steps")

    plt.plot(x_current,y_current,color=charcoal_color,linewidth=5, label = "Current State")
 #   plt.plot(x_future,y_future,color='red',linewidth=5, label="Future State")

    for i, word in enumerate(words):
       plt.text(x_w[i], y_w[i], word, ha='center', va='center', fontname='Avenir Next', fontweight='light',color=charcoal_color,size=12)
    legend = plt.legend(prop={'family':'Avenir Next'})
    for text in legend.get_texts():
       text.set_color('#425369')
    plt.xticks([])
    plt.yticks([])
# Set the aspect ratio to equal
    plt.axis('equal')
    
    plt.annotate("This figure highlights the greatest gaps, \
                 \nrepresenting the degree of shift required, \
                 \nbetween the average current state (depicted \
                 \nby the thick, dark line) and the potential \
                 \nfuture state (thin,outermost line) across \
                 \nprocesses, technology, talent, measurement,\
                 \nand reporting.",
             xy=(1, 0),  # point to label
             xytext=(3, 0),  # location of text
             ha='left', va='center',
    #         arrowprops=dict(facecolor='black', shrink=0.05),
    #         bbox=dict(boxstyle="round,pad=0.3", fc="whitesmoke", alpha=0.8),
             fontsize=20,
             fontname='Avenir Next', color='#425369'
             )
    

# Show the plot
    plt.savefig(dir_+company+'heptagon_plot.png',dpi=300, bbox_inches='tight')

#%% 
##------------------------------------------------
    import pandas as pd
    import matplotlib.pyplot as plt
#import seaborn as sns 
    import numpy as np 


#%%
    fig2,ax = plt.subplots(1, figsize =(20, 9))

    data = pd.read_csv(dir_+'Strategy_'+company+'_.csv',sep=',')
    strategy = data['max'][:-1]
    dim_st   = data['Dimension'][:-1]
    data = pd.read_csv(dir_+'Talent_'+company+'_.csv',sep=',')
    talent  = data['max'][:-1]
    dim_ta  = data['Dimension'][:-1]
    data = pd.read_csv(dir_+'Processes_'+company+'_.csv',sep=',')
    processes  = data['max'][:-1]
    dim_pr  = data['Dimension'][:-1]
    data = pd.read_csv(dir_+'Data_'+company+'_.csv',sep=',')
    data_  = data['max'][:-1]
    dim_da  = data['Dimension'][:-1]
    data = pd.read_csv(dir_+'Technology_'+company+'_.csv',sep=',')
    tech  = data['max'][:-1]
    dim_te  = data['Dimension'][:-1]
    data = pd.read_csv(dir_+'Reporting_'+company+'_.csv',sep=',')
    reporting  = data['max'][:-1]
    dim_re  = data['Dimension'][:-1]
    data = pd.read_csv(dir_+'Measurement_'+company+'_.csv',sep=',')
    measurement  = data['max'][:-1]
    dim_me  = data['Dimension'][:-1]


#%%

##x- and y- axes
    x_ax = (['STRATEGY', 'TALENT', 'PROCESSES', 'DATA', 'MEASUREMENT', 'REPORTING','TECHNOLOGY'])
    y_ax = (['5','4','3','2','1','0'])

    ax.set_xticks([0,1, 2, 3, 4, 5, 6])

    ax.set_xticklabels(x_ax, position=(0,-0.01), font='Avenir Next',color='#425369')
    #ax.set_xlabel('Capability Assessment Dimension',fontsize=25,\
    #                        font='Avenir Next',color='#425369')
    ax.xaxis.labelpad = 20
    ax.tick_params(axis='x', labelsize=17)
    ax.set_xlim(-0.5, 6.5)


    ax.set_yticks([1,2,3,4,5,6])
    ax.set_yticklabels(y_ax,color='#425369')
    ax.tick_params(axis='y', labelsize=15.5)
   # ax.set_ylabel('Steps Between Current State and Future State',fontsize=16,\
   #           font='Avenir Next',color='#425369')
    ax.yaxis.labelpad = 18
    ax.set_ylim(0.5, 6.5)

    c_s = ['#425369']
    c_ta= ['#425369']
    c_p = ['#425369']
    c_d = ['#425369']
    c_m = ['#425369']
    c_r = ['#425369']
    c_te= ['#425369']


##bars and lines
    plt.bar(x_ax, [6,6,6,6,6,6,6], align='center', alpha=0.3,color='grey')
    plt.grid(False)
    st_wor = []
    plt.hlines(6, 0-0.4, 0+0.4,linewidth=45, color='darkgrey')
    plt.hlines(6, 1-0.4, 1+0.4,linewidth=45, color='darkgrey')
    plt.hlines(6, 2-0.4, 2+0.4,linewidth=45, color='darkgrey')
    plt.hlines(6, 3-0.4, 3+0.4,linewidth=45, color='darkgrey')
    plt.hlines(6, 4-0.4, 4+0.4,linewidth=45, color='darkgrey')
    plt.hlines(6, 5-0.4, 5+0.4,linewidth=45, color='darkgrey')
    plt.hlines(6, 6-0.4, 6+0.4,linewidth=45, color='darkgrey')
    
    
    for i in range(len(dim_st)):
       plt.hlines(strategy[i]+1, 0-0.4, 0+0.4,linewidth=55, color=c_s)
      # plt.text(-0.01 * (i+1), strategy[i]+1.3-0.1*i, dim_st.iloc[i], ha='center', va='top', color='whitesmoke', fontsize=9,fontweight='bold')
#     plt.text(-0.4,5-strategy[i],dim_st[i],fontsize=10,font='Avenir Next',color='whitesmoke')
    for i in range(len(dim_ta)):
       plt.hlines(talent[i]+1, 1-0.4, 1+0.4,linewidth=55, color=c_ta)
       #plt.text(0.99+0 * (i+1), talent[i]+1.3-0.15*i, dim_ta.iloc[i], ha='center', va='top', color='whitesmoke', fontsize=9,fontweight='bold')

#       plt.text(0.7+0.1 * (i+1), talent[i]+1, str(i+1), ha='center', va='bottom', color='whitesmoke', fontsize=15)
#    plt.text(0.6,5-talent[i],dim_ta[i],fontsize=10,font='Avenir Next',color='whitesmoke')
    for i in range(len(dim_pr)):
       plt.hlines(processes[i]+1, 2.0-0.4, 2.0+0.4,linewidth=55, color=c_p)
      # plt.text(1.99+0 * (i+1), processes[i]+1.3-0.15*i, dim_pr.iloc[i], ha='center', va='top', color='whitesmoke', fontsize=9,fontweight='bold')

#       plt.text(1.7+0.1 * (i+1), processes[i]+1, str(i+1), ha='center', va='bottom', color='whitesmoke', fontsize=15)
#    plt.text(1.6,5-processes[i],dim_pr[i],fontsize=10,font='Avenir Next',color='whitesmoke')   
    for i in range(len(dim_da)):
       plt.hlines(data_[i]+1, 3-0.4, 3+0.4,linewidth=55, color=c_d)
      # plt.text(2.99+0 * (i+1), data_[i]+1.3-0.15*i, dim_da.iloc[i], ha='center', va='top', color='whitesmoke', fontsize=9,fontweight='bold')

 #      plt.text(2.7+0.1 * (i+1), data_[i]+1, str(i+1), ha='center', va='bottom', color='whitesmoke', fontsize=15)
#    plt.text(2.6,5-data_[i],dim_da[i],fontsize=10,font='Avenir Next',color='whitesmoke')   
    for i in range(len(dim_me)):
       plt.hlines(measurement[i]+1, 4-0.4, 4+0.4,linewidth=55, color=c_m)
      # plt.text(3.99, measurement[i]+1.3-0.15*i, dim_me.iloc[i], ha='center', va='top', color='whitesmoke', fontsize=9, fontweight='bold')
#       plt.text(3.7+0.1 * (i+1), measurement[i]+1, str(i+1), ha='center', va='bottom', color='whitesmoke', fontsize=15)
#    plt.text(3.6,5-measurement[i],dim_me[i],fontsize=10,font='Avenir Next',color='whitesmoke')   
    for i in range(len(dim_re)):
       plt.hlines(reporting[i]+1, 5-0.4, 5+0.4,linewidth=55, color=c_r)
      # plt.text(4.99+0 * (i+1), reporting[i]+1.3-0.15*i, dim_re.iloc[i], ha='center', va='top', color='whitesmoke', fontsize=9,fontweight='bold')

#       plt.text(4.7+0.1 * (i+1), reporting[i]+1, str(i+1), ha='center', va='bottom', color='whitesmoke', fontsize=15)
#    plt.text(4.6,5-reporting[i],dim_re[i],fontsize=10,font='Avenir Next',color='whitesmoke')   
    for i in range(len(dim_te)):
       plt.hlines(tech[i]+1, 6-0.4, 6+0.4,linewidth=55, color=c_te)
      # plt.text(5.99+0 * (i+1), tech[i]+1.3-0.15*i, dim_te.iloc[i], ha='center', va='top', color='whitesmoke', fontsize=9, fontweight='bold')
    
    
    # Adjusting text positioning for overlapping texts on hlines
    def adjust_text_position(text_positions, current_position, direction=-1):
        """Adjust text position to avoid overlap, moving up or down."""
        adjusted_position = current_position
        while adjusted_position in text_positions:
            adjusted_position += 0.19 * direction  # Adjust by 0.15 units up or down
        text_positions.add(adjusted_position)
        return adjusted_position

    # Strategy dimension text positioning
    strategy_text_positions = set()
    for i in range(len(dim_st)):
        text_y = strategy[i] + 1.28  # Initial Y position for text
        text_y = adjust_text_position(strategy_text_positions, text_y)
        plt.text(0, text_y, dim_st.iloc[i], ha='center', va='top', color='whitesmoke', fontsize=9, fontweight='bold')

    # Talent dimension text positioning
    talent_text_positions = set()
    for i in range(len(dim_ta)):
        text_y = talent[i] + 1.28  # Initial Y position for text
        text_y = adjust_text_position(talent_text_positions, text_y)
        plt.text(1, text_y, dim_ta.iloc[i], ha='center', va='top', color='whitesmoke', fontsize=9, fontweight='bold')

    # Processes dimension text positioning
    processes_text_positions = set()
    for i in range(len(dim_pr)):
        text_y = processes[i] + 1.28  # Initial Y position for text
        text_y = adjust_text_position(processes_text_positions, text_y)
        plt.text(2, text_y, dim_pr.iloc[i], ha='center', va='top', color='whitesmoke', fontsize=9, fontweight='bold')

    # Data dimension text positioning
    data_text_positions = set()
    for i in range(len(dim_da)):
        text_y = data_[i] + 1.28  # Initial Y position for text
        text_y = adjust_text_position(data_text_positions, text_y)
        plt.text(3, text_y, dim_da.iloc[i], ha='center', va='top', color='whitesmoke', fontsize=9, fontweight='bold')

    # Measurement dimension text positioning
    measurement_text_positions = set()
    for i in range(len(dim_me)):
        text_y = measurement[i] + 1.28  # Initial Y position for text
        text_y = adjust_text_position(measurement_text_positions, text_y)
        plt.text(4, text_y, dim_me.iloc[i], ha='center', va='top', color='whitesmoke', fontsize=9, fontweight='bold')

    # Reporting dimension text positioning
    reporting_text_positions = set()
    for i in range(len(dim_re)):
        text_y = reporting[i] + 1.28  # Initial Y position for text
        text_y = adjust_text_position(reporting_text_positions, text_y)
        plt.text(5, text_y, dim_re.iloc[i], ha='center', va='top', color='whitesmoke', fontsize=9, fontweight='bold')

    # Technology dimension text positioning
    tech_text_positions = set()
    for i in range(len(dim_te)):
        text_y = tech[i] + 1.28 # Initial Y position for text
        text_y = adjust_text_position(tech_text_positions, text_y)
        plt.text(6, text_y, dim_te.iloc[i], ha='center', va='top', color='whitesmoke', fontsize=9, fontweight='bold')

  #     plt.text(5.7+0.1 * (i+1), tech[i]+1, str(i+1), ha='center', va='bottom', color='whitesmoke', fontsize=15)
#    plt.text(5.6,5-tech[i],dim_te[i],fontsize=10,font='Avenir Next',color='whitesmoke')   
    """
# Strategy dimension
# Initialize a dictionary to keep track of the last y-coordinate used for each hline
# Strategy dimension
    last_hline_y = -1
    text_count = 0
    for i in range(len(dim_st)):
      current_hline_y = strategy[i]+1
      if current_hline_y != last_hline_y:
        last_hline_y = current_hline_y
        text_y = current_hline_y + 0.15
        text_count = 0
      else:
        text_count += 1
        text_y = current_hline_y + 0.15 - 0.15 * text_count
      plt.hlines(current_hline_y, 0-0.4, 0+0.4, linewidth=55, color=c_s)
      plt.text(-0.01 * (i+1), text_y, dim_st.iloc[i], ha='center', va='top', color='whitesmoke', fontsize=9, fontweight='bold')
# Talent dimension
    last_hline_y = -1
    last_text_y = -1
    for i in range(len(dim_ta)):
       current_hline_y = talent[i]+1
       if current_hline_y != last_hline_y:
          last_hline_y = current_hline_y
          text_y = current_hline_y + 0.15
          last_text_y = text_y
       else:
          text_y = last_text_y - 0.15
          last_text_y = text_y
          plt.hlines(current_hline_y, 1-0.4, 1+0.4, linewidth=55, color=c_ta)
          plt.text(1, text_y-0.15, dim_ta.iloc[i], ha='center', va='top', color='whitesmoke', fontsize=9, fontweight='bold')
    """
    ##add test to hlines
    plt.text(-0.2,6,'Future State',fontsize=11,font='Avenir Next',color='whitesmoke')
    plt.text(0.8,6,'Future State',fontsize=11,font='Avenir Next',color='whitesmoke')
    plt.text(1.8,6,'Future State',fontsize=11,font='Avenir Next',color='whitesmoke')
    plt.text(2.8,6,'Future State',fontsize=11,font='Avenir Next',color='whitesmoke')
    plt.text(3.8,6,'Future State',fontsize=11,font='Avenir Next',color='whitesmoke')
    plt.text(4.8,6,'Future State',fontsize=11,font='Avenir Next',color='whitesmoke')
    plt.text(5.8,6,'Future State',fontsize=11,font='Avenir Next',color='whitesmoke')
    """
    plt.text(-0.5,6.7, "This figure provides a more detailed view of the gaps (reflected as the number of steps) between current and future state in each element \
                  \nof the Impact Management Capability dimensions. The elements that are driving the greater degree of shift required between current and future \
                  \nstates (in the next year) of the strategy and processes dimensions include purpose alignment and responsibility framework, respectively.​", 
             fontsize=16,
             fontname='Avenir Next',  color='#425369'
             )
    """
    x_limits = plt.xlim()
    y_limits = plt.ylim()

    print("X limits:", x_limits)
    print("Y limits:", y_limits)

    ax.grid(False)
    plt.savefig(dir_+'CapAss.png',dpi=300,bbox_inches='tight')


    import matplotlib.backends.backend_pdf as pdf
    ##Individual plots\
    """
    fig3,ax = plt.subplots(1, figsize =(17, 10))

    data = pd.read_csv(dir_+'Strategy_'+company+'_.csv',sep=',')
    strategy = data['max'][:-1]
    dim_st   = data['Dimension'][:-1]

    print(strategy)
    print(dim_st)
    # Plotting the bar graph for Strategy dimension
    maturity_dict = {'nascent': 1, 'emerging': 2, 'expanding': 3, 'optimising': 4, 'mature': 5}

    charcoal_color = '#425369'
    font = {'family': 'Avenir Next'}
    plt.rc('font', **font, weight='light')

    ax.bar(dim_st, strategy, color=charcoal_color, edgecolor='black', linewidth=1, alpha=0.7)
    for num in strategy:
        for key, value in maturity_dict.items():
           for i, level in enumerate(strategy):
              ax.text(dim_st[i], level + 0.1, list(maturity_dict.keys())[list(maturity_dict.values()).index(level)], color=charcoal_color, fontweight='light', ha='center',fontsize=18)
    plt.xticks(fontsize=25, color=charcoal_color)
    plt.yticks(fontsize=25, color=charcoal_color)
    
    plt.ylim(0, 5)
    plt.tight_layout()
    # Modernizing the graph with seaborn style
    import seaborn as sns
    sns.set_style("whitegrid")
    sns.despine(left=True)

    plt.gca().spines['top'].set_visible(False)
    plt.gca().spines['right'].set_visible(False)
    plt.gca().spines['bottom'].set_color('#dddddd')
    plt.gca().spines['left'].set_color('#dddddd')
    plt.gca().xaxis.label.set_color('#333333')
    plt.gca().yaxis.label.set_color('#333333')
    plt.gca().tick_params(colors='#333333')
      
   # plt.show(bbox_inches='tight')
    plt.savefig('/Users/ehlke/Desktop/strategy.png',dpi=300, bbox_inches='tight')
    
    from pptx import Presentation
    from pptx.util import Inches
    prs = Presentation('/Users/ehlke/Desktop/Client X Capability Assessment Report Template.pptx')  

    # Select the slide where you want to add the image
    slide = prs.slides[11]  # slide_index is the index of the slide
    # Define the width, height, and coordinates for the image
    left = Inches(1.1)
    top = Inches(1.8)
    width = Inches(5.5)
    height = Inches(3.5)
    
    slide.shapes.add_picture('/Users/ehlke/Desktop/strategy.png', left, top, width, height)

    fig4,ax = plt.subplots(1, figsize =(17, 10))

    data = pd.read_csv(dir_+'Talent_'+company+'_.csv',sep=',')
    talent  = data['max'][:-1]
    dim_ta  = data['Dimension'][:-1]

    # Plotting the bar graph for Strategy dimension
    maturity_dict = {'nascent': 1, 'emerging': 2, 'expanding': 3, 'optimising': 4, 'mature': 5}

    charcoal_color = '#425369'
    font = {'family': 'Avenir Next'}
    plt.rc('font', **font, weight='light')

    ax.bar(dim_ta, talent, color=charcoal_color, edgecolor='black', linewidth=1, alpha=0.7)
    for num in talent:
        for key, value in maturity_dict.items():
           for i, level in enumerate(talent):
              ax.text(dim_ta[i], level + 0.1, list(maturity_dict.keys())[list(maturity_dict.values()).index(level)], color=charcoal_color, fontweight='light', ha='center',fontsize=18)
    plt.xticks(fontsize=25, color=charcoal_color)
    plt.yticks(fontsize=25, color=charcoal_color)
    
    plt.ylim(0, 5)
    plt.tight_layout()
    # Modernizing the graph with seaborn style
    import seaborn as sns
    sns.set_style("whitegrid")
    sns.despine(left=True)

    plt.gca().spines['top'].set_visible(False)
    plt.gca().spines['right'].set_visible(False)
    plt.gca().spines['bottom'].set_color('#dddddd')
    plt.gca().spines['left'].set_color('#dddddd')
    plt.gca().xaxis.label.set_color('#333333')
    plt.gca().yaxis.label.set_color('#333333')
    plt.gca().tick_params(colors='#333333')
      
   # plt.show(bbox_inches='tight')
    plt.savefig('/Users/ehlke/Desktop/talent.png',dpi=300, bbox_inches='tight')
    """

from pptx import Presentation
from pptx.util import Inches
import seaborn as sns
import matplotlib.pyplot as plt
import pandas as pd

def create_plot_and_add_to_ppt(category, company, dir_, prs, slide_index):
    fig, ax = plt.subplots(1, figsize =(17, 10))
    plt.rcParams['font.sans-serif'] = 'DejaVu Sans'

    data = pd.read_csv(f'{dir_}{category}_{company}_.csv', sep=',')
    values = data['max'][:-1]
    dimensions = data['Dimension'][:-1]

    maturity_dict = {'Nascent': 1, 'Emerging': 2, 'Expanding': 3, 'Optimising': 4, 'Mature': 5}
    charcoal_color = '#425369'
    font = {'family': 'Dejavu Sans'}
    plt.rc('font', **font, weight='light')

    ax.bar(dimensions, values, color=charcoal_color, edgecolor='black', linewidth=1, alpha=0.7)
    for num in values:
        for key, value in maturity_dict.items():
            for i, level in enumerate(values):
                ax.text(dimensions[i], level + 0.1, list(maturity_dict.keys())[list(maturity_dict.values()).index(level)],color=charcoal_color, fontweight='light', ha='center',fontsize=18)
    plt.xticks(color=charcoal_color, font='Dejavu Sans')
    plt.yticks(color=charcoal_color, font='Dejavu Sans')
    ax.tick_params(axis='both', which='major', labelsize=21)
    ax.tick_params(axis='x', colors=charcoal_color)
    
    plt.ylim(0, 5)
    plt.tight_layout()

    sns.set_style("whitegrid")
    sns.despine(left=True)
    plt.grid(False)
    plt.gca().spines['top'].set_visible(False)
    plt.gca().spines['right'].set_visible(False)
    plt.gca().spines['bottom'].set_color('#dddddd')
    plt.gca().spines['left'].set_color('#dddddd')
    plt.gca().xaxis.label.set_color('#333333')
    plt.gca().yaxis.label.set_color('#333333')
    plt.gca().tick_params(colors='#333333')
      
    img_path = f'/Users/ehlke/Desktop/{category}.png'
    plt.savefig(img_path, dpi=300, bbox_inches='tight')

    # Add the image to the slide
    slide = prs.slides[slide_index]
    left = Inches(0.9)
    top = Inches(1.8)
    width = Inches(5.8)
    height = Inches(3.5)
    slide.shapes.add_picture(img_path, left, top, width, height)

# Load the presentation
prs = Presentation('/Users/ehlke/Desktop/Client X Capability Assessment Report Template.pptx')  
company='Quantum Foundation'
date_ = 'January 2024'
relativ = 'Relativ Impact'

from pptx.util import Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

slide = prs.slides[0]
txBox = slide.shapes.add_textbox(Inches(9), Inches(0.5), Inches(2), Inches(0.5))
tf = txBox.text_frame
tf.text = company
tf.text = company.upper()

p = tf.paragraphs[0]
p.font.bold = False
p.font.size = Pt(28)
p.font.color.rgb = RGBColor(66, 83, 105)
p.alignment = PP_ALIGN.RIGHT
p.font.name = 'Lato'

# Add Relativ to the slide
slide = prs.slides[0]
txBox_date = slide.shapes.add_textbox(Inches(9), Inches(4), Inches(2), Inches(0.5))
tf_rel = txBox_date.text_frame
tf_rel.text = relativ
tf_rel.text = relativ.upper()

p_rel = tf_rel.paragraphs[0]
p_rel.font.bold = True
p_rel.font.size = Pt(21)
p_rel.font.color.rgb = RGBColor(66, 83, 105)
p_rel.alignment = PP_ALIGN.RIGHT
p_rel.font.name = 'Lato'

# Add date to the slide
slide = prs.slides[0]
txBox_date = slide.shapes.add_textbox(Inches(9), Inches(4.4), Inches(2), Inches(0.5))
tf_date = txBox_date.text_frame
tf_date.text = date_
tf_date.text = date_.upper()

p_date = tf_date.paragraphs[0]
p_date.font.bold = False
p_date.font.size = Pt(18)
p_date.font.color.rgb = RGBColor(66, 83, 105)
p_date.alignment = PP_ALIGN.RIGHT
p_date.font.name = 'Lato'


# Set the font to be super light for all text in the slides
for slide in prs.slides:
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.weight = 'light'





# Call the function for each category
dir_ = '/Users/ehlke/Desktop/Capability_Assessment_Results/'
company='Quantum Foundation'
categories = ['Strategy', 'Talent', 'Processes', 'Data', 'Measurement', 'Reporting', 'Technology']
slide_indices = [11, 15, 20, 24, 28, 33, 37]  # Add more slide indices as needed
for category, slide_index in zip(categories, slide_indices):
    create_plot_and_add_to_ppt(category, company, dir_, prs, slide_index)

# Save the presentation
slide= prs.slides[6]
left = Inches(1.1)
top = Inches(2.1)
width = Inches(8.6)
height = Inches(3.5)

slide.shapes.add_picture(dir_+company+'heptagon_plot.png', left, top, width, height)

slide= prs.slides[7]
left = Inches(1.6)
top = Inches(1.8)
width = Inches(8.2)
height = Inches(3.65)

slide.shapes.add_picture(dir_+'CapAss.png', left, top, width, height)

 
prs.save('/Users/ehlke/Desktop/test.pptx')


capability_report(company='Quantum Foundation')